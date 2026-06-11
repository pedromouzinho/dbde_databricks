# =============================================================================
# pptx_engine.py — DBDE PowerPoint generation engine v3.0 (Gold Standard)
# =============================================================================
# Generates premium PPTX presentations with dark/light contrast, visual cards,
# native charts with data labels, and modern executive styling.
# Supports 12 slide types including native charts, comparisons, process flows,
# and stat+chart combos. Speaker notes on every slide.
#
# Brand guidelines:
#   - Primary font: Montserrat (ExtraBold titles, Bold subtitles, Regular body)
#   - Secondary font: Trebuchet MS (fallback)
#   - Brand accent: #D1005D (cerise)
#   - Dark background: #1A1A2E (midnight navy)
#   - Dark text: #585857 (gray)
#   - Slide size: 13.333 x 7.5 inches (widescreen 16:9)
#   - Dark cover/closing/section slides for visual impact
#   - KPI cards with background fills
#   - Native charts with data labels: bar, pie, line, doughnut
#   - All components fully editable for downstream customisation
# =============================================================================

import io
import logging
import re
import unicodedata
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------
# Palette aligned to the official Millennium template style guide
# (parte_1 slides "PALETA"/"CORES"/"TÍTULO"): cerise accent, navy/charcoal darks,
# a blue family for charts, and Montserrat (titles) + Trebuchet MS (body).
BRAND_ACCENT_HEX = "D1005D"        # cerise — primary accent (official)
BRAND_ACCENT_DARK_HEX = "D40462"   # darker cerise (official variant)
BRAND_DARK_TEXT_HEX = "595959"     # brand gray for body text
BRAND_LIGHT_BG_HEX = "F2F2F2"
BRAND_WHITE_HEX = "FFFFFF"
BRAND_BLACK_HEX = "1A1A1A"
BRAND_HEADER_BADGE_TEXT = "DBDE"

# Dark palette for cover/section/closing slides — official navy & charcoal
BRAND_DARK_BG_HEX = "1B365D"       # navy (the "selected" primary dark)
BRAND_DARK_BG_ALT_HEX = "2E3641"   # charcoal
BRAND_ACCENT_LIGHT_HEX = "F2D0E0"  # soft cerise tint
BRAND_CARD_BG_HEX = "F7F7FA"       # off-white card background
BRAND_CARD_BORDER_HEX = "E8E8EE"   # subtle card border
# Chart/series family — official brand blues (no teal/amber: off-brand)
BRAND_BLUE_HEX = "2F5EA3"          # brand blue (mid)
BRAND_TEAL_HEX = "5C8BD0"          # brand blue (light-mid)
BRAND_AMBER_HEX = "9BB8D3"         # brand blue (light)
BRAND_LEFT_TINT_HEX = "E7EEF5"     # comparison left (official light blue tint)
BRAND_RIGHT_TINT_HEX = "FDF0F4"    # comparison right (warm cerise tint)

# Slide dimensions in EMU (English Metric Units — 914400 EMU = 1 inch)
SLIDE_WIDTH_EMU = 12192000   # 13.333 inches
SLIDE_HEIGHT_EMU = 6858000   # 7.5 inches

# Font sizes in pt
TITLE_FONT_SIZE = 36
SUBTITLE_FONT_SIZE = 18
BODY_FONT_SIZE = 12
SMALL_FONT_SIZE = 10
SECTION_NUMBER_FONT_SIZE = 72
KPI_NUMBER_FONT_SIZE = 54
TABLE_HEADER_FONT_SIZE = 10
TABLE_BODY_FONT_SIZE = 9
BADGE_FONT_SIZE = 8

# Font families
FONT_PRIMARY = "Montserrat"
FONT_SECONDARY = "Trebuchet MS"

# Layout margins (EMU)
MARGIN_LEFT = 609600     # ~0.67 inches
MARGIN_TOP = 914400      # 1 inch
MARGIN_RIGHT = 609600
MARGIN_BOTTOM = 457200   # 0.5 inches
CONTENT_WIDTH = SLIDE_WIDTH_EMU - MARGIN_LEFT - MARGIN_RIGHT


def _emu(inches: float) -> int:
    """Convert inches to EMU."""
    return int(inches * 914400)


def _pt(points: int):
    """Convert points to Pt object."""
    from pptx.util import Pt
    return Pt(points)


def _rgb(hex_color: str):
    """Convert hex color string to RGBColor."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_font(run, *, size: int = BODY_FONT_SIZE, bold: bool = False,
              italic: bool = False, color: str = BRAND_DARK_TEXT_HEX,
              font_name: str = FONT_PRIMARY):
    """Apply font styling to a run."""
    from pptx.util import Pt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    run.font.name = font_name


def _add_bullet_formatting(paragraph, level: int = 0):
    """Add proper bullet character, color and indentation via OOXML XML.

    Level 0: brand-colored bullet '•' with hanging indent
    Level 1: gray en-dash '–' with deeper indent
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    pPr = paragraph._p.get_or_add_pPr()

    # Remove any existing bullet settings
    for tag_suffix in ('buNone', 'buChar', 'buAutoNum', 'buFont', 'buClr', 'buSzPct'):
        for child in list(pPr):
            if child.tag.endswith(tag_suffix):
                pPr.remove(child)

    # Bullet size relative to text (100% = same size)
    buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
    buSzPct.set('val', '100000')  # 100%

    # Bullet font — use Arial for reliable bullet rendering
    buFont = etree.SubElement(pPr, qn('a:buFont'))
    buFont.set('typeface', 'Arial')
    buFont.set('panose', '020B0604020202020204')

    # Bullet color
    buClr = etree.SubElement(pPr, qn('a:buClr'))
    srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
    srgbClr.set('val', BRAND_ACCENT_HEX if level == 0 else BRAND_DARK_TEXT_HEX)

    # Bullet character
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', '\u2022' if level == 0 else '\u2013')  # • or –

    # Indentation: marL = left margin, indent = hanging indent (negative)
    if level == 0:
        pPr.set('marL', '457200')    # 0.5 inch margin
        pPr.set('indent', '-228600')  # 0.25 inch hanging indent
    else:
        pPr.set('marL', '914400')    # 1.0 inch margin
        pPr.set('indent', '-228600')  # 0.25 inch hanging indent

    # Line spacing: 1.2x for readability
    lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.set('val', '120000')  # 120%


def _add_text_box(slide, left, top, width, height, text: str, *,
                  size: int = BODY_FONT_SIZE, bold: bool = False,
                  color: str = BRAND_DARK_TEXT_HEX, alignment=None,
                  font_name: str = FONT_PRIMARY, word_wrap: bool = True):
    """Add a text box with styled text to a slide."""
    from pptx.util import Pt, Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = str(text)
    if alignment:
        p.alignment = alignment
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = str(text)
    _set_font(run, size=size, bold=bold, color=color, font_name=font_name)
    return txBox


def _add_footer_branding(slide, label: str = BRAND_HEADER_BADGE_TEXT):
    """No-op — footer branding removed for cleaner slides.

    Previously added a thin cerise line + small label. Now slides are
    completely clean, relying only on consistent typography and color
    palette for brand identity.
    """
    # Intentionally empty — no footer elements added
    pass


# Legacy reference kept for compatibility — previously contained:
#     bar = slide.shapes.add_shape(1, 0, line_top, SLIDE_WIDTH_EMU, line_h)
#     _add_text_box(slide, ..., label, size=7, bold=True, alignment=PP_ALIGN.RIGHT)


def _add_speaker_notes(slide, text: str):
    """Add speaker notes to a slide."""
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = str(text)[:2000]


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, hex_color: str):
    """Set solid background color on a slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _add_shape_rect(slide, left, top, width, height, fill_hex: str, *, line_hex: str = ""):
    """Add a filled rectangle shape (card, divider, accent bar)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _rgb(line_hex)
        from pptx.util import Pt
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _build_title_slide(prs, title: str, subtitle: str = "",
                       badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium dark cover slide with cerise accent."""
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Dark background
    _set_slide_bg(slide, BRAND_DARK_BG_HEX)

    # Left cerise accent stripe (vertical)
    _add_shape_rect(slide, 0, 0, _emu(0.1), SLIDE_HEIGHT_EMU, BRAND_ACCENT_HEX)

    # Subtle bottom cerise line
    _add_shape_rect(slide, 0, SLIDE_HEIGHT_EMU - _emu(0.04), SLIDE_WIDTH_EMU, _emu(0.04), BRAND_ACCENT_HEX)

    # Title — vertically centred, white on dark
    _add_text_box(
        slide, MARGIN_LEFT + _emu(0.5), _emu(2.0), CONTENT_WIDTH - _emu(0.5), _emu(1.5),
        title, size=44, bold=True, color=BRAND_WHITE_HEX,
        font_name=FONT_PRIMARY,
    )

    # Subtitle — light gray
    if subtitle:
        _add_text_box(
            slide, MARGIN_LEFT + _emu(0.5), _emu(3.6), CONTENT_WIDTH - _emu(0.5), _emu(0.8),
            subtitle, size=SUBTITLE_FONT_SIZE, bold=False,
            color=BRAND_ACCENT_LIGHT_HEX,
        )

    # Date bottom-left, muted
    date_str = datetime.now(timezone.utc).strftime("%d/%m/%Y")
    _add_text_box(
        slide, MARGIN_LEFT + _emu(0.5), _emu(6.3), _emu(3), _emu(0.3),
        date_str, size=SMALL_FONT_SIZE, color="8888AA",
    )

    # Badge label bottom-right
    if badge_text:
        _add_text_box(
            slide, SLIDE_WIDTH_EMU - _emu(2.5), _emu(6.3), _emu(2), _emu(0.3),
            badge_text, size=BADGE_FONT_SIZE, bold=True, color="8888AA",
            alignment=PP_ALIGN.RIGHT,
        )

    _add_footer_branding(slide, badge_text)
    return slide


def _build_section_divider(prs, section_number: int, section_title: str,
                           badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium section divider with dark background and cerise accent."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Dark navy background
    _set_slide_bg(slide, BRAND_DARK_BG_ALT_HEX)

    # Cerise accent bar on left
    _add_shape_rect(slide, 0, 0, _emu(0.1), SLIDE_HEIGHT_EMU, BRAND_ACCENT_HEX)

    # Large section number in cerise
    num_text = f"{section_number:02d}" if section_number < 100 else str(section_number)
    _add_text_box(
        slide, MARGIN_LEFT + _emu(0.5), _emu(1.5), _emu(2.5), _emu(2.0),
        num_text, size=SECTION_NUMBER_FONT_SIZE, bold=True,
        color=BRAND_ACCENT_HEX, font_name=FONT_PRIMARY,
    )

    # Section title — white on dark
    _add_text_box(
        slide, MARGIN_LEFT + _emu(3.5), _emu(2.2), _emu(8), _emu(1.2),
        section_title, size=TITLE_FONT_SIZE, bold=True,
        color=BRAND_WHITE_HEX,
    )

    # Thin cerise underline below title
    _add_shape_rect(
        slide, MARGIN_LEFT + _emu(3.5), _emu(3.5), _emu(2.0), _emu(0.04),
        BRAND_ACCENT_HEX,
    )

    _add_footer_branding(slide, badge_text)
    return slide


def _build_content_slide(prs, title: str, bullets: List[str],
                         badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a content slide with title, cerise accent, and bullet points."""
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title — 28pt for strong visual hierarchy over 12pt body
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.65), CONTENT_WIDTH, _emu(0.6),
        title, size=28, bold=True, color=BRAND_BLACK_HEX,
    )
    # Cerise accent bar under title
    _add_shape_rect(
        slide, MARGIN_LEFT, _emu(1.3), _emu(1.5), _emu(0.04), BRAND_ACCENT_HEX,
    )

    # Bullets
    if bullets:
        txBox = slide.shapes.add_textbox(
            MARGIN_LEFT, _emu(1.7), CONTENT_WIDTH, _emu(5.0),
        )
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True

        for idx, bullet_text in enumerate(bullets):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Support sub-bullets with "- " or "• " prefix
            text = str(bullet_text).strip()
            is_sub = text.startswith("- ") or text.startswith("• ")
            if is_sub:
                text = text[2:].strip()
                level = 1
            else:
                level = 0

            p.level = level
            p.space_after = Pt(8)
            if idx == 0:
                p.space_before = Pt(4)

            # Apply proper bullet formatting via XML
            _add_bullet_formatting(p, level=level)

            run = p.add_run()
            run.text = text
            font_size = BODY_FONT_SIZE if level == 0 else (BODY_FONT_SIZE - 1)
            _set_font(run, size=font_size, color=BRAND_DARK_TEXT_HEX)

    _add_footer_branding(slide, badge_text)
    return slide


def _build_two_column_slide(prs, title: str, left_content: List[str],
                            right_content: List[str],
                            badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a two-column content slide with subtle divider."""
    from pptx.util import Pt, Emu
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.65), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )
    # Cerise accent bar under title
    _add_shape_rect(slide, MARGIN_LEFT, _emu(1.3), _emu(1.5), _emu(0.04), BRAND_ACCENT_HEX)

    col_width = (CONTENT_WIDTH - _emu(0.5)) // 2

    # Subtle vertical divider
    divider_x = MARGIN_LEFT + col_width + _emu(0.2)
    _add_shape_rect(slide, divider_x, _emu(1.6), _emu(0.02), _emu(5.0), BRAND_CARD_BORDER_HEX)

    for col_idx, items in enumerate([left_content, right_content]):
        col_left = MARGIN_LEFT if col_idx == 0 else (MARGIN_LEFT + col_width + _emu(0.5))
        if not items:
            continue
        txBox = slide.shapes.add_textbox(col_left, _emu(1.7), col_width, _emu(5.0))
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True
        for idx, item_text in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            text = str(item_text).strip()
            is_sub = text.startswith("- ") or text.startswith("• ")
            if is_sub:
                text = text[2:].strip()
                level = 1
            else:
                level = 0
            p.level = level
            p.space_after = Pt(8)
            _add_bullet_formatting(p, level=level)
            run = p.add_run()
            run.text = text
            font_size = BODY_FONT_SIZE if level == 0 else (BODY_FONT_SIZE - 1)
            _set_font(run, size=font_size, color=BRAND_DARK_TEXT_HEX)

    _add_footer_branding(slide, badge_text)
    return slide


def _build_kpi_slide(prs, title: str, kpis: List[Dict[str, Any]],
                     badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium KPI/metrics slide with visual cards.

    Each KPI dict: {value: str, label: str, description?: str}
    Max 4 KPIs per slide.
    """
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.65), CONTENT_WIDTH, _emu(0.6),
        title, size=28, bold=True, color=BRAND_BLACK_HEX,
    )

    n = min(len(kpis), 4)
    if n == 0:
        return slide

    # Card layout with gaps
    card_gap = _emu(0.3)
    total_gap = card_gap * (n - 1)
    card_width = (CONTENT_WIDTH - total_gap) // n
    card_height = _emu(3.5)
    card_top = _emu(1.7)

    for i, kpi in enumerate(kpis[:4]):
        x = MARGIN_LEFT + i * (card_width + card_gap)

        # Card background
        _add_shape_rect(slide, x, card_top, card_width, card_height,
                        BRAND_CARD_BG_HEX, line_hex=BRAND_CARD_BORDER_HEX)

        # Cerise top accent bar on card
        _add_shape_rect(slide, x, card_top, card_width, _emu(0.06), BRAND_ACCENT_HEX)

        # KPI value (large number)
        _add_text_box(
            slide, x + _emu(0.2), card_top + _emu(0.5),
            card_width - _emu(0.4), _emu(1.2),
            str(kpi.get("value", "")),
            size=KPI_NUMBER_FONT_SIZE, bold=True, color=BRAND_ACCENT_HEX,
            alignment=PP_ALIGN.CENTER,
        )
        # KPI label
        _add_text_box(
            slide, x + _emu(0.2), card_top + _emu(1.8),
            card_width - _emu(0.4), _emu(0.5),
            str(kpi.get("label", "")),
            size=BODY_FONT_SIZE, bold=True, color=BRAND_DARK_TEXT_HEX,
            alignment=PP_ALIGN.CENTER,
        )
        # KPI description (optional)
        desc = kpi.get("description", "")
        if desc:
            _add_text_box(
                slide, x + _emu(0.2), card_top + _emu(2.4),
                card_width - _emu(0.4), _emu(0.8),
                str(desc),
                size=SMALL_FONT_SIZE, bold=False, color=BRAND_DARK_TEXT_HEX,
                alignment=PP_ALIGN.CENTER,
            )

    _add_footer_branding(slide, badge_text)
    return slide


def _build_table_slide(prs, title: str, headers: List[str],
                       rows: List[List[str]],
                       badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium slide with a branded data table."""
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.65), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )
    # Cerise accent bar under title
    _add_shape_rect(slide, MARGIN_LEFT, _emu(1.3), _emu(1.5), _emu(0.04), BRAND_ACCENT_HEX)

    if not headers:
        return slide

    max_rows = min(len(rows), 15)  # cap at 15 rows per slide
    n_cols = len(headers)
    n_rows = max_rows + 1  # +1 for header

    table_top = _emu(1.7)
    table_height = _emu(min(5.0, 0.35 * n_rows + 0.1))
    col_width = CONTENT_WIDTH // n_cols

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        MARGIN_LEFT, table_top, CONTENT_WIDTH, table_height,
    )
    table = table_shape.table

    # Style header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(header)
        # Dark navy header background (premium)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(BRAND_DARK_BG_HEX)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, size=TABLE_HEADER_FONT_SIZE, bold=True,
                          color=BRAND_WHITE_HEX)

    # Style data rows with alternating colors
    for ri, row_data in enumerate(rows[:max_rows]):
        for ci in range(n_cols):
            cell = table.cell(ri + 1, ci)
            cell_val = row_data[ci] if ci < len(row_data) else ""
            cell.text = str(cell_val)
            # Zebra striping
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(BRAND_WHITE_HEX)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(BRAND_LIGHT_BG_HEX)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_font(run, size=TABLE_BODY_FONT_SIZE,
                              color=BRAND_DARK_TEXT_HEX)

    if len(rows) > max_rows:
        note = f"(+{len(rows) - max_rows} linhas omitidas)"
        _add_text_box(
            slide, MARGIN_LEFT, table_top + table_height + _emu(0.15),
            CONTENT_WIDTH, _emu(0.3),
            note, size=8, italic=True, color=BRAND_DARK_TEXT_HEX,
        )

    _add_footer_branding(slide, badge_text)
    return slide


def _build_agenda_slide(prs, items: List[str],
                        badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create an agenda/index slide."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.85), CONTENT_WIDTH, _emu(0.6),
        "Agenda", size=TITLE_FONT_SIZE, bold=True, color=BRAND_BLACK_HEX,
    )

    for idx, item_text in enumerate(items[:12]):
        y = _emu(1.8 + idx * 0.45)
        # Number circle
        _add_text_box(
            slide, MARGIN_LEFT, y, _emu(0.5), _emu(0.35),
            f"{idx + 1:02d}", size=14, bold=True, color=BRAND_ACCENT_HEX,
            alignment=PP_ALIGN.CENTER,
        )
        # Item text
        _add_text_box(
            slide, MARGIN_LEFT + _emu(0.7), y, _emu(10), _emu(0.35),
            str(item_text), size=14, bold=False, color=BRAND_DARK_TEXT_HEX,
        )

    _add_footer_branding(slide, badge_text)
    return slide


def _build_closing_slide(prs, text: str = "Obrigado",
                         subtitle: str = "",
                         badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium dark closing slide matching the cover."""
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background matching cover
    _set_slide_bg(slide, BRAND_DARK_BG_HEX)

    # Bottom cerise accent line
    _add_shape_rect(slide, 0, SLIDE_HEIGHT_EMU - _emu(0.04), SLIDE_WIDTH_EMU, _emu(0.04), BRAND_ACCENT_HEX)

    # Main text — cerise on dark
    _add_text_box(
        slide, MARGIN_LEFT, _emu(2.2), CONTENT_WIDTH, _emu(1.5),
        text, size=44, bold=True, color=BRAND_ACCENT_HEX,
        alignment=PP_ALIGN.CENTER,
    )

    if subtitle:
        _add_text_box(
            slide, MARGIN_LEFT, _emu(3.9), CONTENT_WIDTH, _emu(0.8),
            subtitle, size=SUBTITLE_FONT_SIZE, bold=False,
            color=BRAND_ACCENT_LIGHT_HEX, alignment=PP_ALIGN.CENTER,
        )

    date_str = datetime.now(timezone.utc).strftime("%d/%m/%Y")
    _add_text_box(
        slide, MARGIN_LEFT, _emu(6.3), CONTENT_WIDTH, _emu(0.3),
        date_str, size=SMALL_FONT_SIZE, color="8888AA",
        alignment=PP_ALIGN.CENTER,
    )

    _add_footer_branding(slide, badge_text)
    return slide


# ---------------------------------------------------------------------------
# New slide builders — chart, comparison, process, stat_chart
# ---------------------------------------------------------------------------

def _style_chart_premium(chart, series_list, chart_type: str):
    """Apply premium styling to a chart: data labels, brand colors, clean grid."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    # Apply brand colors
    _brand_chart_colors = [
        BRAND_ACCENT_HEX, BRAND_BLUE_HEX, BRAND_DARK_TEXT_HEX,
        BRAND_AMBER_HEX, BRAND_TEAL_HEX,
    ]
    plot = chart.plots[0]
    for idx, s in enumerate(plot.series):
        color_hex = _brand_chart_colors[idx % len(_brand_chart_colors)]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = _rgb(color_hex)

    # Data labels — show values on all non-pie charts
    if chart_type not in ("pie", "doughnut"):
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_value = True
        data_labels.show_category_name = False
        data_labels.show_series_name = False
        data_labels.font.size = Pt(8)
        data_labels.font.name = FONT_PRIMARY
        data_labels.font.color.rgb = _rgb(BRAND_DARK_TEXT_HEX)
    else:
        # Pie/doughnut: show percentages
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_value = False
        data_labels.show_category_name = True
        data_labels.font.size = Pt(8)
        data_labels.font.name = FONT_PRIMARY

    # Legend
    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT_PRIMARY

    # Axis styling (pie/doughnut have no axes)
    try:
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(9)
        cat_axis.tick_labels.font.name = FONT_PRIMARY
        cat_axis.tick_labels.font.color.rgb = _rgb(BRAND_DARK_TEXT_HEX)
        cat_axis.format.line.fill.background()  # hide axis line
        cat_axis.has_major_gridlines = False
    except (ValueError, AttributeError):
        pass
    try:
        val_axis = chart.value_axis
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.name = FONT_PRIMARY
        val_axis.tick_labels.font.color.rgb = _rgb(BRAND_DARK_TEXT_HEX)
        val_axis.format.line.fill.background()  # hide axis line
        # Subtle gridlines
        val_axis.has_major_gridlines = True
        val_axis.major_gridlines.format.line.color.rgb = _rgb("E2E8F0")
        val_axis.major_gridlines.format.line.width = Pt(0.5)
    except (ValueError, AttributeError):
        pass


def _build_chart_slide(prs, title: str, chart_type: str,
                       categories: List[str], series: List[Dict[str, Any]],
                       badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium chart slide with data labels and clean styling.

    chart_type: "bar", "column", "line", "pie", "doughnut"
    categories: ["Q1", "Q2", "Q3", ...]
    series: [{"name": "Revenue", "values": [100, 200, 300]}, ...]
    """
    from pptx.util import Pt, Emu
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.4), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    # Map chart type string to XL_CHART_TYPE
    _chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    # Pie/doughnut only support 1 series — auto-upgrade to column if multiple
    if chart_type in ("pie", "doughnut") and len(series) > 1:
        chart_type = "column"
        logger.info("[PptxEngine] %s chart with %d series → upgraded to column",
                    chart_type, len(series))

    xl_type = _chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Build chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories[:12]
    for s in series[:5]:
        name = str(s.get("name", ""))
        values = s.get("values", [])
        numeric_values = []
        for v in values[:12]:
            try:
                numeric_values.append(float(v))
            except (ValueError, TypeError):
                numeric_values.append(0)
        chart_data.add_series(name, numeric_values)

    # Add chart to slide
    chart_left = MARGIN_LEFT + _emu(0.3)
    chart_top = _emu(1.2)
    chart_width = CONTENT_WIDTH - _emu(0.6)
    chart_height = _emu(5.3)
    chart_frame = slide.shapes.add_chart(
        xl_type, chart_left, chart_top, chart_width, chart_height, chart_data,
    )

    # Premium chart styling
    _style_chart_premium(chart_frame.chart, series, chart_type)

    _add_footer_branding(slide, badge_text)
    return slide


def _build_comparison_slide(prs, title: str,
                            left_title: str, left_items: List[str],
                            right_title: str, right_items: List[str],
                            badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium comparison slide with tinted column backgrounds.

    Two columns with distinct colored backgrounds and headers.
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.4), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    col_gap = _emu(0.3)
    col_width = (CONTENT_WIDTH - col_gap) // 2
    col_top = _emu(1.15)
    col_height = _emu(5.7)
    right_x = MARGIN_LEFT + col_width + col_gap

    # Column background cards with tints
    _add_shape_rect(slide, MARGIN_LEFT, col_top, col_width, col_height,
                    BRAND_LEFT_TINT_HEX, line_hex=BRAND_CARD_BORDER_HEX)
    _add_shape_rect(slide, right_x, col_top, col_width, col_height,
                    BRAND_RIGHT_TINT_HEX, line_hex=BRAND_CARD_BORDER_HEX)

    # Column headers with cerise background
    header_h = _emu(0.5)
    _add_shape_rect(slide, MARGIN_LEFT, col_top, col_width, header_h, BRAND_BLUE_HEX)
    _add_text_box(
        slide, MARGIN_LEFT + _emu(0.2), col_top + _emu(0.05),
        col_width - _emu(0.4), _emu(0.4),
        left_title, size=14, bold=True, color=BRAND_WHITE_HEX,
    )

    _add_shape_rect(slide, right_x, col_top, col_width, header_h, BRAND_ACCENT_HEX)
    _add_text_box(
        slide, right_x + _emu(0.2), col_top + _emu(0.05),
        col_width - _emu(0.4), _emu(0.4),
        right_title, size=14, bold=True, color=BRAND_WHITE_HEX,
    )

    # Column items
    items_top = col_top + header_h + _emu(0.15)
    for col_items, col_x in [(left_items, MARGIN_LEFT), (right_items, right_x)]:
        if not col_items:
            continue
        txBox = slide.shapes.add_textbox(
            col_x + _emu(0.15), items_top,
            col_width - _emu(0.3), col_height - header_h - _emu(0.3),
        )
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(col_items[:8]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            _add_bullet_formatting(p, level=0)
            run = p.add_run()
            run.text = str(item).strip()
            _set_font(run, size=11, color=BRAND_DARK_TEXT_HEX)

    _add_footer_branding(slide, badge_text)
    return slide


def _build_process_slide(prs, title: str, steps: List[Dict[str, Any]],
                         badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium horizontal process/timeline slide with cards.

    steps: [{"label": "Step 1", "description": "Details"}, ...]
    Max 5 steps for readability.
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.4), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    n = min(len(steps), 5)
    if n == 0:
        _add_footer_branding(slide, badge_text)
        return slide

    step_gap = _emu(0.15)
    step_width = (CONTENT_WIDTH - step_gap * (n - 1)) // n
    circle_size = _emu(0.9)
    circle_y = _emu(1.8)

    for i, step in enumerate(steps[:5]):
        step_x = MARGIN_LEFT + i * (step_width + step_gap)
        cx = step_x + (step_width - circle_size) // 2

        # Card background below the circle
        card_top = circle_y + circle_size + _emu(0.15)
        card_height = _emu(3.8)
        _add_shape_rect(slide, step_x, card_top, step_width, card_height,
                        BRAND_CARD_BG_HEX, line_hex=BRAND_CARD_BORDER_HEX)
        # Cerise top edge on card
        _add_shape_rect(slide, step_x, card_top, step_width, _emu(0.04), BRAND_ACCENT_HEX)

        # Numbered circle
        circle = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            cx, circle_y, circle_size, circle_size,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _rgb(BRAND_ACCENT_HEX)
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        _set_font(run, size=20, bold=True, color=BRAND_WHITE_HEX)

        # Arrow connector between circles (chevron style)
        if i < n - 1:
            arrow_left = cx + circle_size + _emu(0.05)
            arrow_width = (step_width + step_gap) - circle_size - _emu(0.1)
            arrow_y = circle_y + circle_size // 2 - _emu(0.1)
            # Use right arrow shape (MSO_SHAPE.NOTCHED_RIGHT_ARROW = 94)
            try:
                arrow = slide.shapes.add_shape(
                    94, arrow_left, arrow_y, arrow_width, _emu(0.2),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = _rgb(BRAND_ACCENT_HEX)
                arrow.line.fill.background()
            except Exception:
                # Fallback to simple line
                connector = slide.shapes.add_shape(
                    1, arrow_left, arrow_y + _emu(0.08), arrow_width, _emu(0.04),
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = _rgb(BRAND_ACCENT_HEX)
                connector.line.fill.background()

        # Step label inside card
        label = str(step.get("label", f"Passo {i + 1}"))
        _add_text_box(
            slide, step_x + _emu(0.1), card_top + _emu(0.2),
            step_width - _emu(0.2), _emu(0.5),
            label, size=12, bold=True, color=BRAND_BLACK_HEX,
            alignment=PP_ALIGN.CENTER,
        )

        # Step description inside card
        desc = str(step.get("description", ""))
        if desc:
            _add_text_box(
                slide, step_x + _emu(0.1), card_top + _emu(0.75),
                step_width - _emu(0.2), _emu(2.8),
                desc, size=10, bold=False, color=BRAND_DARK_TEXT_HEX,
                alignment=PP_ALIGN.CENTER,
            )

    _add_footer_branding(slide, badge_text)
    return slide


def _build_stat_chart_slide(prs, title: str,
                            stat_value: str, stat_label: str,
                            chart_type: str, categories: List[str],
                            series: List[Dict[str, Any]],
                            badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a premium combo slide: KPI card on left, chart on right.

    Ideal for "headline number + supporting trend" layouts.
    """
    from pptx.util import Pt, Emu
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.4), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    # Left side: KPI card
    stat_x = MARGIN_LEFT
    stat_width = _emu(3.8)
    stat_top = _emu(1.5)
    stat_height = _emu(4.0)

    # Card background
    _add_shape_rect(slide, stat_x, stat_top, stat_width, stat_height,
                    BRAND_CARD_BG_HEX, line_hex=BRAND_CARD_BORDER_HEX)
    # Cerise top bar
    _add_shape_rect(slide, stat_x, stat_top, stat_width, _emu(0.06), BRAND_ACCENT_HEX)

    _add_text_box(
        slide, stat_x + _emu(0.2), stat_top + _emu(0.8), stat_width - _emu(0.4), _emu(1.5),
        str(stat_value), size=54, bold=True, color=BRAND_ACCENT_HEX,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, stat_x + _emu(0.2), stat_top + _emu(2.3), stat_width - _emu(0.4), _emu(0.6),
        str(stat_label), size=14, bold=True, color=BRAND_DARK_TEXT_HEX,
        alignment=PP_ALIGN.CENTER,
    )

    # Right side: chart
    _chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    if chart_type in ("pie", "doughnut") and len(series) > 1:
        chart_type = "column"
    xl_type = _chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = categories[:12]
    for s in series[:3]:
        name = str(s.get("name", ""))
        values = []
        for v in s.get("values", [])[:12]:
            try:
                values.append(float(v))
            except (ValueError, TypeError):
                values.append(0)
        chart_data.add_series(name, values)

    chart_left = MARGIN_LEFT + _emu(4.3)
    chart_width = CONTENT_WIDTH - _emu(4.3)
    chart_frame = slide.shapes.add_chart(
        xl_type, chart_left, _emu(1.3), chart_width, _emu(5.2), chart_data,
    )

    # Premium chart styling
    _style_chart_premium(chart_frame.chart, series, chart_type)

    _add_footer_branding(slide, badge_text)
    return slide


# ---------------------------------------------------------------------------
# Smart slide validation & auto-correction
# ---------------------------------------------------------------------------
# This layer enforces presentation quality rules REGARDLESS of what the LLM
# sends. It splits overloaded slides, trims excess content, and ensures
# professional structure.
# ---------------------------------------------------------------------------

_MAX_BULLETS_PER_SLIDE = 7
_MAX_BULLET_LENGTH = 150
_MAX_TITLE_LENGTH = 80
_MAX_KPIS_PER_SLIDE = 4
_MAX_TABLE_ROWS_PER_SLIDE = 12
_MAX_TABLE_COLS = 8
_MAX_AGENDA_ITEMS = 12
_MAX_TWO_COL_ITEMS = 8
_MAX_CHART_CATEGORIES = 12
_MAX_CHART_SERIES = 5
_MAX_PROCESS_STEPS = 5
_MAX_COMPARISON_ITEMS = 8
_MAX_NOTES_LENGTH = 2000


def _validate_and_fix_slides(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Validate and auto-correct slide specs for professional quality.

    Rules enforced:
    1. Content slides with >7 bullets → split into multiple slides
    2. Bullet text >150 chars → truncate with ellipsis
    3. Title >80 chars → truncate
    4. KPI slides with >4 KPIs → split
    5. Table with >12 rows → split into continuation slides
    6. Table with >8 columns → trim to 8
    7. Two-column slides with >8 items per side → trim
    8. Empty slides → removed
    9. Adjacent content slides with same title → merge if under bullet limit
    10. Consecutive sections without content between → drop duplicate
    """
    if not slides:
        return []

    fixed = []
    for spec in slides:
        if not isinstance(spec, dict):
            continue
        slide_type = str(spec.get("type", "content")).lower().strip()

        # --- Truncate title ---
        title = str(spec.get("title", "")).strip()
        if len(title) > _MAX_TITLE_LENGTH:
            title = title[:_MAX_TITLE_LENGTH - 1].rstrip() + "…"
            spec = {**spec, "title": title}

        # --- Content slides: split if too many bullets ---
        if slide_type in ("content", "bullets"):
            bullets = spec.get("bullets", [])
            if isinstance(bullets, str):
                bullets = [b.strip() for b in bullets.split("\n") if b.strip()]
            # Truncate individual bullets
            bullets = [
                (b[:_MAX_BULLET_LENGTH - 1].rstrip() + "…" if len(b) > _MAX_BULLET_LENGTH else b)
                for b in bullets
            ]
            if not bullets:
                # Empty content slide → skip
                continue
            # Split into chunks of _MAX_BULLETS_PER_SLIDE
            for chunk_idx in range(0, len(bullets), _MAX_BULLETS_PER_SLIDE):
                chunk = bullets[chunk_idx:chunk_idx + _MAX_BULLETS_PER_SLIDE]
                chunk_title = title
                if chunk_idx > 0:
                    chunk_title = f"{title} (cont.)"
                fixed.append({**spec, "title": chunk_title, "bullets": chunk})
            continue

        # --- KPI slides: split if >4 ---
        elif slide_type in ("kpi", "kpis", "metrics"):
            kpis = spec.get("kpis", [])
            if not kpis:
                continue
            for chunk_idx in range(0, len(kpis), _MAX_KPIS_PER_SLIDE):
                chunk = kpis[chunk_idx:chunk_idx + _MAX_KPIS_PER_SLIDE]
                chunk_title = title
                if chunk_idx > 0:
                    chunk_title = f"{title} (cont.)"
                fixed.append({**spec, "title": chunk_title, "kpis": chunk})
            continue

        # --- Table slides: split rows, trim columns ---
        elif slide_type == "table":
            headers = spec.get("headers", [])
            rows = spec.get("rows", [])
            if not headers:
                continue
            # Trim columns
            if len(headers) > _MAX_TABLE_COLS:
                headers = headers[:_MAX_TABLE_COLS]
                rows = [r[:_MAX_TABLE_COLS] for r in rows]
            if not rows:
                fixed.append({**spec, "headers": headers, "rows": []})
                continue
            # Split rows
            for chunk_idx in range(0, len(rows), _MAX_TABLE_ROWS_PER_SLIDE):
                chunk = rows[chunk_idx:chunk_idx + _MAX_TABLE_ROWS_PER_SLIDE]
                chunk_title = title
                if chunk_idx > 0:
                    chunk_title = f"{title} (cont.)"
                fixed.append({
                    **spec, "title": chunk_title,
                    "headers": headers, "rows": chunk,
                })
            continue

        # --- Two-column: trim items ---
        elif slide_type in ("two_column", "two_columns"):
            left = spec.get("left", [])
            right = spec.get("right", [])
            if not left and not right:
                continue
            fixed.append({
                **spec,
                "left": left[:_MAX_TWO_COL_ITEMS],
                "right": right[:_MAX_TWO_COL_ITEMS],
            })
            continue

        # --- Agenda: trim items ---
        elif slide_type in ("agenda", "index"):
            items = spec.get("items", [])
            if not items:
                continue
            fixed.append({**spec, "items": items[:_MAX_AGENDA_ITEMS]})
            continue

        # --- Chart: validate categories and series ---
        elif slide_type == "chart":
            categories = spec.get("categories", [])
            series = spec.get("series", [])
            if not categories or not series:
                continue
            fixed.append({
                **spec,
                "categories": categories[:_MAX_CHART_CATEGORIES],
                "series": series[:_MAX_CHART_SERIES],
            })
            continue

        # --- Comparison: trim items ---
        elif slide_type == "comparison":
            left_items = spec.get("left_items", [])
            right_items = spec.get("right_items", [])
            if not left_items and not right_items:
                continue
            fixed.append({
                **spec,
                "left_items": left_items[:_MAX_COMPARISON_ITEMS],
                "right_items": right_items[:_MAX_COMPARISON_ITEMS],
            })
            continue

        # --- Process: trim steps ---
        elif slide_type in ("process", "timeline"):
            steps = spec.get("steps", [])
            if not steps:
                continue
            fixed.append({**spec, "steps": steps[:_MAX_PROCESS_STEPS]})
            continue

        # --- Stat+chart: validate ---
        elif slide_type in ("stat_chart", "stat"):
            # Degrade to kpi if no chart data
            categories = spec.get("categories", [])
            series = spec.get("series", [])
            if not categories or not series:
                degraded = {
                    "type": "kpi",
                    "title": title,
                    "kpis": [{"value": spec.get("stat_value", ""),
                              "label": spec.get("stat_label", "")}],
                }
                if spec.get("notes"):
                    degraded["notes"] = spec["notes"]
                fixed.append(degraded)
            else:
                fixed.append({
                    **spec,
                    "categories": categories[:_MAX_CHART_CATEGORIES],
                    "series": series[:_MAX_CHART_SERIES],
                })
            continue

        # --- All other types pass through ---
        else:
            fixed.append(spec)

    # --- Post-pass: remove consecutive empty sections ---
    cleaned = []
    for i, slide in enumerate(fixed):
        slide_type = str(slide.get("type", "")).lower().strip()
        if slide_type in ("section", "section_divider", "divider"):
            # Check if next slide is also a section → skip this one
            if i + 1 < len(fixed):
                next_type = str(fixed[i + 1].get("type", "")).lower().strip()
                if next_type in ("section", "section_divider", "divider"):
                    continue
        cleaned.append(slide)

    return cleaned


_SLIDE_TYPE_MAP = {
    "title": _build_title_slide,
    "cover": _build_title_slide,
    "capa": _build_title_slide,
    "section": _build_section_divider,
    "section_divider": _build_section_divider,
    "divider": _build_section_divider,
    "content": _build_content_slide,
    "bullets": _build_content_slide,
    "two_column": _build_two_column_slide,
    "two_columns": _build_two_column_slide,
    "kpi": _build_kpi_slide,
    "kpis": _build_kpi_slide,
    "metrics": _build_kpi_slide,
    "table": _build_table_slide,
    "agenda": _build_agenda_slide,
    "index": _build_agenda_slide,
    "chart": _build_chart_slide,
    "comparison": _build_comparison_slide,
    "process": _build_process_slide,
    "timeline": _build_process_slide,
    "stat_chart": _build_stat_chart_slide,
    "stat": _build_stat_chart_slide,
    "closing": _build_closing_slide,
    "end": _build_closing_slide,
    "obrigado": _build_closing_slide,
}


def _build_slide_from_spec(prs, spec: Dict[str, Any], section_counter: int,
                           badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Build a single slide from a spec dict.

    Spec format (12 types):
    {
        "type": "content|title|section|kpi|table|two_column|agenda|closing
                |chart|comparison|process|stat_chart",
        "title": "...",
        "subtitle": "...",         # for title/closing
        "bullets": ["...", ...],   # for content
        "left": ["...", ...],      # for two_column
        "right": ["...", ...],     # for two_column
        "kpis": [{value, label, description?}, ...],  # for kpi
        "headers": ["...", ...],   # for table
        "rows": [["...", ...], ...], # for table
        "items": ["...", ...],     # for agenda
        "section_number": int,     # for section (auto-calculated if omitted)
        "text": "...",             # for closing
        "chart_type": "bar|column|line|pie|doughnut",  # for chart/stat_chart
        "categories": ["...", ...],  # for chart/stat_chart
        "series": [{"name": "...", "values": [...]}, ...],  # for chart/stat_chart
        "left_title": "...",       # for comparison
        "left_items": ["...", ...],  # for comparison
        "right_title": "...",      # for comparison
        "right_items": ["...", ...],  # for comparison
        "steps": [{"label": "...", "description": "..."}, ...],  # for process
        "stat_value": "...",       # for stat_chart
        "stat_label": "...",       # for stat_chart
        "notes": "...",            # speaker notes (all types)
    }
    """
    slide_type = str(spec.get("type", "content")).lower().strip()
    builder = _SLIDE_TYPE_MAP.get(slide_type, _build_content_slide)

    title = str(spec.get("title", "")).strip()
    subtitle = str(spec.get("subtitle", "")).strip()
    notes = str(spec.get("notes", "")).strip()

    slide = None
    try:
        if slide_type in ("title", "cover", "capa"):
            slide = builder(prs, title or "Apresentação", subtitle, badge_text)

        elif slide_type in ("section", "section_divider", "divider"):
            num = spec.get("section_number", section_counter)
            slide = builder(prs, int(num), title, badge_text)

        elif slide_type in ("content", "bullets"):
            bullets = spec.get("bullets", [])
            if isinstance(bullets, str):
                bullets = [b.strip() for b in bullets.split("\n") if b.strip()]
            slide = builder(prs, title, bullets, badge_text)

        elif slide_type in ("two_column", "two_columns"):
            left = spec.get("left", [])
            right = spec.get("right", [])
            slide = builder(prs, title, left, right, badge_text)

        elif slide_type in ("kpi", "kpis", "metrics"):
            kpis = spec.get("kpis", [])
            slide = builder(prs, title, kpis, badge_text)

        elif slide_type in ("table",):
            headers = spec.get("headers", [])
            rows = spec.get("rows", [])
            slide = builder(prs, title, headers, rows, badge_text)

        elif slide_type in ("agenda", "index"):
            items = spec.get("items", [])
            slide = builder(prs, items, badge_text)

        elif slide_type in ("chart",):
            chart_type = str(spec.get("chart_type", "column"))
            categories = spec.get("categories", [])
            series = spec.get("series", [])
            slide = builder(prs, title, chart_type, categories, series, badge_text)

        elif slide_type in ("comparison",):
            left_title = str(spec.get("left_title", ""))
            left_items = spec.get("left_items", [])
            right_title = str(spec.get("right_title", ""))
            right_items = spec.get("right_items", [])
            slide = builder(prs, title, left_title, left_items,
                            right_title, right_items, badge_text)

        elif slide_type in ("process", "timeline"):
            steps = spec.get("steps", [])
            slide = builder(prs, title, steps, badge_text)

        elif slide_type in ("stat_chart", "stat"):
            stat_value = str(spec.get("stat_value", ""))
            stat_label = str(spec.get("stat_label", ""))
            chart_type = str(spec.get("chart_type", "column"))
            categories = spec.get("categories", [])
            series = spec.get("series", [])
            slide = builder(prs, title, stat_value, stat_label,
                            chart_type, categories, series, badge_text)

        elif slide_type in ("closing", "end", "obrigado"):
            text = spec.get("text", title or "Obrigado")
            slide = builder(prs, text, subtitle, badge_text)

        else:
            # Fallback: treat as content slide
            bullets = spec.get("bullets", [])
            if not bullets and spec.get("text"):
                bullets = [spec["text"]]
            slide = _build_content_slide(prs, title, bullets, badge_text)

    except Exception as e:
        logger.warning("[PptxEngine] slide build error (type=%s): %s", slide_type, e)
        slide = _build_content_slide(
            prs, title or "Slide",
            [f"(Erro ao gerar slide: {str(e)[:100]})"],
            badge_text,
        )

    # Apply speaker notes universally
    if slide and notes:
        _add_speaker_notes(slide, notes)

    return slide


# ---------------------------------------------------------------------------
# Brand logo stamping (official Millennium mark)
# ---------------------------------------------------------------------------
_DARK_SLIDE_TYPES = {"title", "cover", "capa", "section", "section_divider",
                     "divider", "closing", "end", "obrigado"}


def _stamp_branding(prs, slide_types) -> None:
    """Stamp the official brand logo on every slide: a small footer mark on
    content slides and a larger mark on cover/closing. Dark slides use the white
    logo variant so the (dark) mark stays visible. No-op if the asset is missing.
    """
    import os
    assets = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
    logo_dark = os.path.join(assets, "brand_logo.png")          # for light slides
    logo_light = os.path.join(assets, "brand_logo_white.png")   # for dark slides
    if not os.path.exists(logo_dark):
        return
    from pptx.util import Emu
    try:
        from PIL import Image
        with Image.open(logo_dark) as im:
            aspect = (im.size[0] / im.size[1]) if im.size[1] else (101 / 87)
    except Exception:
        aspect = 101 / 87
    inch = 914400
    for i, slide in enumerate(prs.slides):
        stype = slide_types[i] if i < len(slide_types) else "content"
        dark_bg = stype in _DARK_SLIDE_TYPES
        logo = logo_light if (dark_bg and os.path.exists(logo_light)) else logo_dark
        prominent = stype in ("title", "cover", "capa", "closing", "end", "obrigado")
        h_emu = int((0.62 if prominent else 0.30) * inch)
        w_emu = int(h_emu * aspect)
        if prominent:
            left = SLIDE_WIDTH_EMU - w_emu - int(0.60 * inch)
            top = int(0.50 * inch)
        else:
            left = SLIDE_WIDTH_EMU - w_emu - int(0.45 * inch)
            top = SLIDE_HEIGHT_EMU - h_emu - int(0.28 * inch)
        try:
            slide.shapes.add_picture(logo, Emu(left), Emu(top), height=Emu(h_emu))
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_presentation(
    title: str,
    slides: List[Dict[str, Any]],
    *,
    subtitle: str = "",
    badge_text: str = BRAND_HEADER_BADGE_TEXT,
    include_title_slide: bool = True,
    include_closing_slide: bool = True,
) -> io.BytesIO:
    """Generate a professional PPTX presentation (12 slide types, native charts).

    Args:
        title: Presentation title
        slides: List of slide spec dicts (see _build_slide_from_spec)
        subtitle: Subtitle for title slide
        badge_text: Footer label (default: "DBDE")
        include_title_slide: Auto-add title slide if not in specs
        include_closing_slide: Auto-add closing slide if not in specs

    Returns:
        BytesIO buffer with PPTX content
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()

    # Set widescreen 16:9 dimensions
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    # ── Smart validation: enforce quality rules before rendering ──
    validated_slides = _validate_and_fix_slides(slides)

    # Check if slides already start with a title/cover
    has_title_slide = False
    has_closing_slide = False
    if validated_slides:
        first_type = str(validated_slides[0].get("type", "")).lower().strip()
        last_type = str(validated_slides[-1].get("type", "")).lower().strip()
        has_title_slide = first_type in ("title", "cover", "capa")
        has_closing_slide = last_type in ("closing", "end", "obrigado")

    built_types: List[str] = []
    # Auto-add title slide (with default notes)
    if include_title_slide and not has_title_slide:
        title_slide = _build_title_slide(prs, title, subtitle, badge_text)
        _add_speaker_notes(title_slide, f"Apresentação: {title}")
        built_types.append("title")

    # Build each slide
    section_counter = 1
    for spec in validated_slides:
        if not isinstance(spec, dict):
            continue
        slide_type = str(spec.get("type", "content")).lower().strip()
        built_types.append(slide_type)
        if slide_type in ("section", "section_divider", "divider"):
            _build_slide_from_spec(prs, spec, section_counter, badge_text)
            section_counter += 1
        else:
            _build_slide_from_spec(prs, spec, section_counter, badge_text)

    # Auto-add closing slide (with default notes)
    if include_closing_slide and not has_closing_slide:
        closing_slide = _build_closing_slide(prs, "Obrigado", "", badge_text)
        _add_speaker_notes(closing_slide, "Agradecer e abrir para questões")
        built_types.append("closing")

    # Stamp the official brand logo on every slide (footer + cover/closing)
    _stamp_branding(prs, built_types)

    # Save to buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_presentation_from_outline(
    title: str,
    outline: str,
    *,
    subtitle: str = "",
    badge_text: str = BRAND_HEADER_BADGE_TEXT,
) -> io.BytesIO:
    """Generate presentation from a plain-text outline.

    Parses a simple outline format:
    - Lines starting with "# " = section dividers
    - Lines starting with "## " = content slide titles
    - Lines starting with "- " under a ## = bullets
    - Lines starting with "| " = table rows (first row = headers)

    Returns BytesIO buffer with PPTX.
    """
    slides = []
    current_slide = None
    section_num = 0

    for raw_line in outline.split("\n"):
        line = raw_line.strip()
        if not line:
            continue

        if line.startswith("# "):
            # Save previous slide
            if current_slide:
                slides.append(current_slide)
                current_slide = None
            section_num += 1
            slides.append({
                "type": "section",
                "title": line[2:].strip(),
                "section_number": section_num,
            })

        elif line.startswith("## "):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "type": "content",
                "title": line[3:].strip(),
                "bullets": [],
            }

        elif line.startswith("- ") or line.startswith("• "):
            if current_slide is None:
                current_slide = {"type": "content", "title": "", "bullets": []}
            if "bullets" not in current_slide:
                current_slide["bullets"] = []
            current_slide["bullets"].append(line[2:].strip())

        elif line.startswith("| "):
            # Table detection
            if current_slide and current_slide.get("type") != "table":
                if current_slide.get("bullets") or current_slide.get("title"):
                    slides.append(current_slide)
                current_slide = {
                    "type": "table",
                    "title": current_slide.get("title", ""),
                    "headers": [],
                    "rows": [],
                }
            if current_slide is None:
                current_slide = {"type": "table", "title": "", "headers": [], "rows": []}
            cells = [c.strip() for c in line.strip("| ").split("|")]
            # Skip separator rows (---)
            if all(re.match(r'^-+$', c) for c in cells):
                continue
            if not current_slide.get("headers"):
                current_slide["headers"] = cells
            else:
                current_slide["rows"].append(cells)

        else:
            # Plain text line — add as bullet to current slide
            if current_slide and "bullets" in current_slide:
                current_slide["bullets"].append(line)

    if current_slide:
        slides.append(current_slide)

    return generate_presentation(
        title, slides,
        subtitle=subtitle,
        badge_text=badge_text,
    )


# ---------------------------------------------------------------------------
# Opus-powered slide planner — AI-driven slide structure generation
# ---------------------------------------------------------------------------
# When the conversation LLM (any tier) calls generate_presentation, it can
# pass a `content` string instead of pre-structured slides. The planner then
# uses Claude Opus 4.6 (pro tier) to produce professional slide specs.
# ---------------------------------------------------------------------------

_SLIDE_PLANNER_PROMPT = """Tu és um consultor sénior da DBDE que cria apresentações de nível board-ready.
Pensas como estratega: a apresentação deve contar uma história, não listar factos.

PRINCÍPIOS DE DESIGN PREMIUM:
- Cada slide deve ter UM propósito claro e UMA mensagem
- O deck tem arco narrativo: contexto → diagnóstico → evidência → implicações → ação
- Dados visuais (charts, KPIs, tables) > bullets sempre que possível
- Nunca 2 slides "content" seguidos — quebra com chart, kpi, comparison ou table
- Títulos assertivos que comunicam a conclusão, não o tema (ex: "Gap total em 11 domínios" > "Análise de gaps")
- Speaker notes contêm o detalhe que o apresentador precisa (contexto, dados, talking points)

TIPOS DE SLIDE DISPONÍVEIS (12 tipos — usa pelo menos 5 diferentes):

1. "title" — Capa escura premium (dark navy + cerise accent)
   Campos: title, subtitle, notes?

2. "section" — Divisor visual (dark background, número grande)
   Campos: title, notes?

3. "content" — Slide com bullets (max 6, curtos)
   Campos: title, bullets (array, max 6), notes?

4. "two_column" — Duas colunas de bullets
   Campos: title, left (array), right (array), notes?

5. "kpi" — Cards visuais com números grandes (max 4)
   Campos: title, kpis (array de {value, label, description?}), notes?
   USAR SEMPRE QUE HOUVER MÉTRICAS — é o slide com maior impacto visual

6. "table" — Tabela com header dark navy
   Campos: title, headers (array), rows (array de arrays), notes?

7. "agenda" — Índice numerado
   Campos: items (array de strings), notes?

8. "chart" — Gráfico nativo com data labels
   Campos: title, chart_type ("bar"|"column"|"line"|"pie"|"doughnut"),
   categories (array), series (array de {name, values}), notes?
   OBRIGATÓRIO: pelo menos 1 chart por deck. Dados numéricos NUNCA vão em bullets.

9. "comparison" — Lado a lado com colunas coloridas (azul vs cerise)
   Campos: title, left_title, left_items (array), right_title, right_items (array), notes?
   USAR para qualquer before/after, prós/contras, entidade A vs B

10. "process" — Timeline horizontal com cards e setas (max 5 passos)
    Campos: title, steps (array de {label, description}), notes?

11. "stat_chart" — Card KPI + gráfico lado a lado
    Campos: title, stat_value, stat_label, chart_type, categories, series, notes?

12. "closing" — Slide final escuro (dark navy matching capa)
    Campos: text, subtitle, notes?

REGRAS OBRIGATÓRIAS:
1. 8-14 slides — focado, sem padding
2. Max 6 bullets por slide, cada um < 100 chars
3. Pelo menos 1 chart ou stat_chart
4. Pelo menos 1 kpi
5. Pelo menos 1 comparison ou table
6. Nunca 2 content slides seguidos
7. Notes em TODOS os slides (detalhado, para o apresentador)
8. Títulos assertivos em português de Portugal
9. Dados numéricos → chart/kpi/table, NUNCA em bullets
10. Arco narrativo: diagnóstico → evidência → ação

Responde APENAS com um JSON array. Sem explicações, sem markdown, só o JSON.
Exemplo:
[
  {"type": "title", "title": "Transformação Digital Empresas", "subtitle": "Análise estratégica e roadmap — Comissão Executiva", "notes": "Enquadrar o objetivo da sessão."},
  {"type": "agenda", "items": ["Diagnóstico macro", "Gap funcional por domínio", "Implicações estratégicas", "Roadmap recomendado"], "notes": "Agenda para 30 minutos de sessão executiva."},
  {"type": "kpi", "title": "Sinais de Alerta", "kpis": [{"value": "0%", "label": "Readiness digital", "description": "Nenhuma capacidade digital coberta"}, {"value": "11", "label": "Domínios com gap", "description": "Totalidade dos domínios analisados"}], "notes": "Mensagem central: gap total."},
  {"type": "chart", "title": "Amplitude Funcional por Domínio", "chart_type": "bar", "categories": ["Pagamentos","Analytics","Conta"], "series": [{"name": "Benchmark", "values": [10,10,9]}, {"name": "Millennium", "values": [0,0,0]}], "notes": "Contrastar amplitude do benchmark com zero cobertura."},
  {"type": "comparison", "title": "Posicionamento Competitivo", "left_title": "Revolut Business", "left_items": ["Conta multimoeda +35 divisas","Cartões virtuais self-service","Analytics tempo real"], "right_title": "Millennium BCP", "right_items": ["Sem oferta multimoeda digital","Emissão dependente de balcão","Sem dashboards analíticos"], "notes": "Slide de impacto — cada linha é um gap direto."},
  {"type": "process", "title": "Roadmap Recomendado", "steps": [{"label": "Core Digital", "description": "Conta multimoeda e pagamentos"},{"label": "Cartões", "description": "Self-service e controlos"},{"label": "Analytics", "description": "Dashboards e reporting"},{"label": "Ecossistema", "description": "APIs e integrações"}], "notes": "4 fases sequenciais por prioridade de impacto."},
  {"type": "closing", "text": "Decisão Recomendada", "subtitle": "Aprovar business case para as Fases 1 e 2", "notes": "Pedir decisão executiva concreta."}
]"""

_BENCHMARK_EXEC_APPENDIX = """
MODO BENCHMARK / COMPARAÇÃO EXECUTIVA:
- Age como consultor McKinsey/BCG: sintetiza, não resume
- 8-12 slides max
- Obrigatório: kpi de síntese + chart de gaps + table de domínios + comparison + process/roadmap
- Abre com diagnóstico macro (números fortes), desenvolve com evidência visual, fecha com ação
- Usa chart com 2 séries para benchmark lado a lado (entidade A vs B)
- Títulos que comunicam a conclusão: "Gap total em todos os domínios" > "Análise de gaps"
- Speaker notes com os dados qualitativos e talking points para o apresentador
- Não repitas informação entre slides; cada slide acrescenta valor novo
- Se houver dados de API/integração, dedica 1 slide de comparison ou table
"""

_NARRATIVE_SHAPER_PROMPT = """Tu és um consultor sénior da DBDE. A tua tarefa é analisar o conteúdo fornecido e criar um briefing narrativo estruturado para uma apresentação.

NÃO cries slides. Cria um briefing que outro sistema vai usar para gerar slides.

O briefing deve conter:
1. CLASSIFICAÇÃO: tipo de deck (benchmark, status_update, strategy, training, operational, pitch)
2. AUDIÊNCIA: executiva, técnica, ou mista
3. MENSAGEM CENTRAL: 1 frase que resume a conclusão principal
4. ARCO NARRATIVO: 3-5 actos da história (ex: "Contexto → Diagnóstico → Evidência → Implicações → Ação")
5. DADOS CHAVE: lista de factos/números/métricas mais relevantes para visualização
6. COMPARAÇÕES: pares de contraste identificados (ex: entidade A vs B, antes vs depois)
7. RECOMENDAÇÃO: o call-to-action final

Responde em JSON com esta estrutura:
{
  "classification": "benchmark",
  "audience": "executive",
  "central_message": "...",
  "narrative_arc": ["Contexto", "Diagnóstico", "Evidência", "Implicações", "Ação"],
  "key_metrics": [{"value": "...", "label": "...", "context": "..."}],
  "comparisons": [{"entity_a": "...", "entity_b": "...", "dimensions": ["..."]}],
  "key_data_points": ["..."],
  "recommendation": "..."
}

Responde APENAS com JSON. Sem explicações."""


def _sanitize_planner_content(content: str) -> str:
    """Remove leaked planner metadata and clean up content for the AI planner.

    Strips patterns like "SLIDE 5 — PAGAMENTOS" that sometimes leak from
    previous LLM planning into visible content. Also removes excessive
    whitespace and normalizes structure markers.
    """
    if not content:
        return ""

    lines = content.split("\n")
    cleaned = []
    for line in lines:
        stripped = line.strip()
        # Remove "SLIDE N — TITLE" or "SLIDE N:" patterns
        if re.match(r"^SLIDE\s+\d+\s*[—–\-:]\s*", stripped, re.IGNORECASE):
            continue
        # Remove "--- Slide N ---" separator patterns
        if re.match(r"^-{3,}\s*slide\s+\d+\s*-{3,}$", stripped, re.IGNORECASE):
            continue
        cleaned.append(line)

    return "\n".join(cleaned)


def _normalize_text_key(text: str) -> str:
    """Lowercase, strip accents and normalize spaces for fuzzy matching."""
    raw = unicodedata.normalize("NFKD", str(text or ""))
    ascii_text = "".join(ch for ch in raw if not unicodedata.combining(ch))
    return re.sub(r"\s+", " ", ascii_text).strip().lower()


def _looks_like_benchmark_content(content: str, title: str = "", context: str = "") -> bool:
    """Detect benchmark/comparison style content with strong PDF-derived signals."""
    text = _normalize_text_key("\n".join([title or "", context or "", content or ""]))
    signals = 0
    if "benchmark" in text or "comparacao" in text:
        signals += 1
    if "situacao millennium bcp" in text:
        signals += 1
    if "impacto para os clientes" in text:
        signals += 1
    if "tabela comparativa" in text:
        signals += 1
    if "revolut business" in text:
        signals += 1
    if "comissao executiva" in text or "comite executivo" in text or "executiva" in text:
        signals += 1
    return signals >= 3


def _looks_like_executive_request(title: str = "", context: str = "", content: str = "") -> bool:
    text = _normalize_text_key("\n".join([title or "", context or "", content or ""]))
    executive_markers = (
        "comissao executiva",
        "comite executivo",
        "executiva",
        "steerco",
        "board",
        "ceo",
        "direcao",
        "decision",
    )
    return any(marker in text for marker in executive_markers)


def _extract_benchmark_sections(content: str) -> List[Dict[str, Any]]:
    """Extract qualitative benchmark domains from PDF-like text.

    Expected pattern:
    Domain heading
    o Feature ...
    Situação Millennium bcp
    • Status ...
    Impacto para os clientes
    • Impact ...
    """
    if not content:
        return []

    qualitative = content.split("2. Tabela comparativa: Revolut Business vs Nova App Millennium bcp")[0]
    raw_lines = [re.sub(r"\s+", " ", line).strip() for line in qualitative.splitlines()]
    lines = [line for line in raw_lines if line]

    sections: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None
    mode = "intro"

    def flush_current():
        nonlocal current
        if not current:
            return
        if current.get("features") or current.get("status") or current.get("impact"):
            sections.append(current)
        current = None

    for idx, line in enumerate(lines):
        next_line = lines[idx + 1] if idx + 1 < len(lines) else ""
        if line.isdigit():
            continue
        if line in {
            "Revolut Business",
            "Benchmark Revolut Business e comparação com o Millennium BCP",
            "1. Lista de funcionalidades da aplicação Revolut Business",
        }:
            continue
        if line.startswith("2. Tabela comparativa"):
            break
        if line == "Situação Millennium bcp":
            mode = "status"
            continue
        if line == "Impacto para os clientes":
            mode = "impact"
            continue
        if line.startswith("o "):
            if current is None:
                continue
            current.setdefault("features", []).append(line[2:].strip())
            mode = "features"
            continue
        if line.startswith("• "):
            if current is None:
                continue
            target_key = "status" if mode == "status" else "impact"
            current.setdefault(target_key, []).append(line[2:].strip())
            continue

        looks_like_heading = (
            not line.startswith(("o ", "• "))
            and next_line.startswith("o ")
            and len(line) <= 80
            and not re.match(r"^\d+\.", line)
        )
        if looks_like_heading:
            flush_current()
            current = {
                "title": line,
                "features": [],
                "status": [],
                "impact": [],
                "narrative": [],
            }
            mode = "features"
            continue

        if current and mode == "impact":
            current.setdefault("narrative", []).append(line)

    flush_current()
    return sections


def _score_benchmark_status(status_bullets: List[str]) -> Tuple[int, int, int]:
    """Estimate current readiness from qualitative Millennium status bullets."""
    if not status_bullets:
        return 0, 0, 0

    strong_positive = (
        "sim",
        "disponivel",
        "disponiveis",
        "mvp",
        "implementacao",
        "visao agrupamento",
        "cobertura",
    )
    mild_positive = ("prevista", "previsto", "fase 2", "segunda fase")
    strong_negative = (
        "nao",
        "ausencia",
        "sem ",
        "manual",
        "nao existe",
        "nao estao",
        "nao disponivel",
        "nao disponiveis",
    )
    mild_negative = ("limitada", "limitado", "apenas", "dependencia", "futuras")

    score = 0
    positive_hits = 0
    negative_hits = 0
    for bullet in status_bullets:
        norm = _normalize_text_key(bullet)
        bullet_score = 0
        if any(marker in norm for marker in strong_positive):
            bullet_score += 2
            positive_hits += 1
        elif any(marker in norm for marker in mild_positive):
            bullet_score += 1
            positive_hits += 1

        if any(marker in norm for marker in strong_negative):
            bullet_score -= 2
            negative_hits += 1
        elif any(marker in norm for marker in mild_negative):
            bullet_score -= 1
            negative_hits += 1
        score += bullet_score

    readiness = max(5, min(95, int(round(50 + (score * 12 / max(1, len(status_bullets)))))))
    return readiness, positive_hits, negative_hits


def _pick_key_gap(status_bullets: List[str], impact_bullets: List[str]) -> Tuple[str, str]:
    """Return a short gap statement and client impact statement."""
    preferred_gap = next(
        (
            bullet for bullet in status_bullets
            if any(
                token in _normalize_text_key(bullet)
                for token in ("nao", "ausencia", "sem ", "manual", "limitad", "nao existe")
            )
        ),
        status_bullets[0] if status_bullets else "Cobertura atual limitada",
    )
    impact = impact_bullets[0] if impact_bullets else "Impacto relevante na proposta de valor empresarial"
    return preferred_gap[:95], impact[:95]


def _build_benchmark_brief(content: str, title: str = "", context: str = "") -> Optional[Dict[str, Any]]:
    """Build an executive benchmark briefing from dense qualitative content."""
    if not _looks_like_benchmark_content(content, title, context):
        return None

    sections = _extract_benchmark_sections(content)
    if len(sections) < 3:
        return None

    domain_rows = []
    section_summaries = []
    comparison_left: List[str] = []
    comparison_right: List[str] = []
    for sec in sections:
        feature_count = len(sec.get("features", []))
        readiness, positive_hits, negative_hits = _score_benchmark_status(sec.get("status", []))
        gap, impact = _pick_key_gap(sec.get("status", []), sec.get("impact", []))
        domain_rows.append({
            "title": sec["title"],
            "feature_count": feature_count,
            "readiness": readiness,
            "positive_hits": positive_hits,
            "negative_hits": negative_hits,
            "gap": gap,
            "impact": impact,
        })
        section_summaries.append({
            "title": sec["title"],
            "features": sec.get("features", [])[:3],
            "status": sec.get("status", [])[:3],
            "impact": sec.get("impact", [])[:2],
            "gap": gap,
        })
        for feature in sec.get("features", [])[:2]:
            if len(comparison_left) < 6:
                comparison_left.append(feature[:90])
        if len(comparison_right) < 6:
            comparison_right.append(gap)

    if not domain_rows:
        return None

    domain_rows.sort(key=lambda row: (row["readiness"], -row["feature_count"]))
    top_gaps = domain_rows[:5]
    avg_readiness = int(round(sum(row["readiness"] for row in domain_rows) / max(1, len(domain_rows))))
    total_features = sum(row["feature_count"] for row in domain_rows)
    critical_domains = sum(1 for row in domain_rows if row["readiness"] <= 30)
    strongest_domain = max(domain_rows, key=lambda row: row["readiness"])
    weakest_domain = min(domain_rows, key=lambda row: row["readiness"])

    executive_points = [
        f"{critical_domains} domínios apresentam gap crítico face ao benchmark",
        f"Readiness média estimada do Millennium: {avg_readiness}%",
        f"Maior vulnerabilidade em {weakest_domain['title']}",
        f"Melhor posicionamento relativo em {strongest_domain['title']}",
    ]

    table_rows = [
        [row["title"], f"{row['readiness']}%", row["gap"], row["impact"]]
        for row in top_gaps
    ]

    briefing_lines = [
        "MODO EXECUTIVO BENCHMARK",
        f"Domínios analisados: {len(domain_rows)}",
        f"Capacidades benchmark identificadas: {total_features}",
        f"Readiness média estimada: {avg_readiness}%",
        f"Domínios com gap crítico: {critical_domains}",
        "",
        "Mensagem executiva:",
    ]
    briefing_lines.extend(f"- {point}" for point in executive_points)
    briefing_lines.extend([
        "",
        "Tabela executiva por domínio:",
        "| Domínio | Feature breadth | Readiness (%) | Gap principal |",
        "| --- | --- | --- | --- |",
    ])
    briefing_lines.extend(
        f"| {row['title']} | {row['feature_count']} | {row['readiness']} | {row['gap']} |"
        for row in domain_rows
    )
    briefing_lines.extend([
        "",
        "Comparação executiva:",
        "Revolut diferencia-se por:",
    ])
    briefing_lines.extend(f"- {item}" for item in comparison_left[:5])
    briefing_lines.append("Millennium precisa reforçar:")
    briefing_lines.extend(f"- {item}" for item in comparison_right[:5])

    return {
        "domain_rows": domain_rows,
        "agenda_items": [
            "Mensagem executiva",
            "Gap funcional por domínio",
            "Top gaps com impacto cliente",
            "Posicionamento competitivo",
            "Implicações e próximos passos",
        ],
        "executive_points": executive_points,
        "table_rows": table_rows,
        "comparison_left": comparison_left[:5],
        "comparison_right": comparison_right[:5],
        "chart_categories": [row["title"] for row in domain_rows[:8]],
        "chart_values": [float(row["readiness"]) for row in domain_rows[:8]],
        "summary_kpis": [
            {"value": str(len(domain_rows)), "label": "Domínios analisados"},
            {"value": f"{avg_readiness}%", "label": "Readiness média"},
            {"value": str(critical_domains), "label": "Gaps críticos"},
            {"value": str(total_features), "label": "Capacidades benchmark"},
        ],
        "sections": section_summaries,
        "briefing_text": "\n".join(briefing_lines),
    }


def _extract_benchmark_api_summary(content: str) -> Optional[Dict[str, Any]]:
    """Extract a compact API benchmark summary from the PDF text."""
    if "3. Tabela comparativa: Revolut API vs Millennium API" not in content:
        return None

    api_part = content.split("3. Tabela comparativa: Revolut API vs Millennium API", 1)[1]
    api_part = api_part.split("Posicionamento da Revolut Business API face ao B2B Millennium API", 1)[0]
    normalized = _normalize_text_key(api_part)
    if "revolut api" not in normalized and "millennium api" not in normalized:
        return None

    left_items = [
        "Integração rápida via REST, sandbox self-service e documentação aberta",
        "Ponto forte em FX, multi-moeda e automação digital",
        "Integrações nativas com ERPs e apps de produtividade",
        "Maior agilidade para embedding financeiro em contextos digitais",
    ]
    right_items = [
        "Cobertura superior em pagamentos ao Estado e serviços",
        "Canal mais robusto de ficheiros, retornos e processos corporativos",
        "Confirmação de beneficiário e profundidade funcional no contexto nacional",
        "Integração mais formal, lenta e com maior esforço de onboarding",
    ]
    table_rows = [
        ["Velocidade de integração", "Forte", "Moderada", "Revolut"],
        ["Pagamentos complexos nacionais", "Limitado", "Forte", "Millennium"],
        ["Ficheiros e retornos", "Parcial", "Forte", "Millennium"],
        ["FX e multi-moeda", "Forte", "Limitado", "Revolut"],
        ["Integrações ERP/apps", "Forte", "Moderada", "Revolut"],
    ]
    return {
        "left_items": left_items,
        "right_items": right_items,
        "table_rows": table_rows,
    }


def _build_benchmark_planner_payload(
    benchmark_brief: Dict[str, Any],
    api_summary: Optional[Dict[str, Any]],
    *,
    title: str = "",
    context: str = "",
) -> str:
    """Build a compact, executive benchmark payload for the slide planner.

    This deliberately favors structured signals over raw PDF text so the LLM
    can plan rich slides without drowning in 20k+ chars of qualitative text.
    """
    lines = [
        f"TÍTULO: {title or 'Benchmark executivo'}",
        f"CONTEXTO: {context or 'Apresentação para comissão executiva'}",
        "OBJETIVO: gerar um deck executivo visual, com comparações, charts e tabelas, não um documento de bullets.",
        "",
        "AGENDA EXECUTIVA:",
    ]
    lines.extend(f"- {item}" for item in benchmark_brief.get("agenda_items", []))
    lines.extend([
        "",
        "KPIS DE ABERTURA:",
    ])
    for kpi in benchmark_brief.get("summary_kpis", [])[:4]:
        lines.append(f"- {kpi.get('label')}: {kpi.get('value')}")

    lines.extend([
        "",
        "DOMÍNIOS E GAPS:",
        "| Domínio | Readiness | Gap principal | Impacto |",
        "| --- | --- | --- | --- |",
    ])
    for row in benchmark_brief.get("table_rows", [])[:6]:
        lines.append(f"| {row[0]} | {row[1]} | {row[2]} | {row[3]} |")

    lines.extend([
        "",
        "POSICIONAMENTO COMPETITIVO:",
        "Revolut reforça:",
    ])
    lines.extend(f"- {item}" for item in benchmark_brief.get("comparison_left", [])[:5])
    lines.append("Millennium precisa reforçar:")
    lines.extend(f"- {item}" for item in benchmark_brief.get("comparison_right", [])[:5])

    if benchmark_brief.get("sections"):
        lines.extend([
            "",
            "DOMÍNIOS PARA SLIDES DE DETALHE:",
        ])
        for sec in benchmark_brief["sections"][:4]:
            lines.append(f"[{sec['title']}]")
            for feat in sec.get("features", [])[:3]:
                lines.append(f"  Revolut/benchmark: {feat}")
            lines.append(f"  Gap Millennium: {sec.get('gap', '')}")
            for impact in sec.get("impact", [])[:2]:
                lines.append(f"  Impacto: {impact}")

    if api_summary:
        lines.extend([
            "",
            "APIs E INTEGRAÇÕES:",
            "Revolut API acelera:",
        ])
        lines.extend(f"- {item}" for item in api_summary.get("left_items", [])[:4])
        lines.append("Millennium API aprofunda:")
        lines.extend(f"- {item}" for item in api_summary.get("right_items", [])[:4])
        lines.extend([
            "Tabela executiva das APIs:",
            "| Tema | Revolut | Millennium | Vantagem |",
            "| --- | --- | --- | --- |",
        ])
        for row in api_summary.get("table_rows", [])[:5]:
            lines.append(f"| {row[0]} | {row[1]} | {row[2]} | {row[3]} |")

    lines.extend([
        "",
        "CONSTRAINTS VISUAIS OBRIGATÓRIOS:",
        "- máximo 12 slides",
        "- pelo menos 2 charts",
        "- pelo menos 2 slides do tipo table/comparison",
        "- evita sections vazias; cada slide deve acrescentar informação",
        "- usar comparação executiva por domínio quando fizer sentido",
    ])
    return "\n".join(lines)


def _parse_json_from_llm(raw: str) -> Optional[Any]:
    """Robustly extract JSON (array or object) from LLM response."""
    import json as _json
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        lines = cleaned.split("\n")
        lines = [l for l in lines if not l.strip().startswith("```")]
        cleaned = "\n".join(lines).strip()
    # Try array first, then object
    for open_ch, close_ch in [("[", "]"), ("{", "}")]:
        start = cleaned.find(open_ch)
        end = cleaned.rfind(close_ch)
        if start != -1 and end > start:
            try:
                return _json.loads(cleaned[start:end + 1])
            except _json.JSONDecodeError:
                continue
    return None


async def _shape_narrative(content: str, title: str = "", context: str = "",
                           *, tier: str = "pro") -> Optional[Dict[str, Any]]:
    """Step 1 of planning: analyze content and produce narrative briefing.

    This shapes the raw content into a structured narrative before slide planning.
    Returns None on failure (caller falls back to direct planning).
    """
    from llm_provider_databricks import llm_simple

    user_parts = []
    if title:
        user_parts.append(f"TÍTULO: {title}")
    if context:
        user_parts.append(f"CONTEXTO: {context[:3000]}")
    user_parts.append(f"CONTEÚDO:\n{content[:15000]}")
    user_prompt = "\n\n".join(user_parts)

    full_prompt = f"{_NARRATIVE_SHAPER_PROMPT}\n\n{user_prompt}"
    try:
        raw = await llm_simple(full_prompt, tier=tier, max_tokens=3000)
        result = _parse_json_from_llm(raw)
        if isinstance(result, dict) and result.get("central_message"):
            logger.info("[PptxPlanner] Narrative shaped: %s / %s",
                        result.get("classification"), result.get("central_message", "")[:60])
            return result
    except Exception as e:
        logger.warning("[PptxPlanner] Narrative shaping failed: %s", e)
    return None


async def plan_slides_with_opus(
    content: str,
    title: str = "",
    context: str = "",
    *,
    tier: str = "pro",
) -> List[Dict[str, Any]]:
    """Two-step AI planning: narrative shaping → slide generation.

    Step 1: Shape narrative (classify, extract key data, define arc)
    Step 2: Plan slides using narrative briefing + content

    Args:
        content: Free-text content/topic for the presentation
        title: Optional title hint
        context: Optional context (conversation history, data, etc.)
        tier: LLM tier to use (default: "pro" = Opus)

    Returns:
        List of slide spec dicts ready for generate_presentation()
    """
    import json as _json
    from llm_provider_databricks import llm_simple

    # Sanitize content: strip leaked planner metadata (e.g. "SLIDE 5 — TITULO")
    sanitized_content = _sanitize_planner_content(content)
    benchmark_brief = _build_benchmark_brief(sanitized_content, title, context)
    executive_mode = _looks_like_executive_request(title, context, sanitized_content) or bool(benchmark_brief)

    api_summary = _extract_benchmark_api_summary(sanitized_content) if benchmark_brief else None
    planner_content = sanitized_content
    if benchmark_brief:
        planner_content = _build_benchmark_planner_payload(
            benchmark_brief,
            api_summary,
            title=title,
            context=context,
        )

    # ── Step 1: Narrative shaping (parallel-safe, non-blocking on failure) ──
    narrative = None
    try:
        narrative = await _shape_narrative(sanitized_content, title, context, tier=tier)
    except Exception as e:
        logger.info("[PptxPlanner] Narrative shaping skipped: %s", e)

    # ── Step 2: Slide planning with enriched context ──
    user_prompt_parts = []
    if title:
        user_prompt_parts.append(f"TÍTULO DA APRESENTAÇÃO: {title}")
    if context:
        user_prompt_parts.append(f"CONTEXTO ADICIONAL:\n{context[:4000]}")

    # Inject narrative briefing if available
    if narrative:
        narrative_text = (
            f"BRIEFING NARRATIVO (usa isto para estruturar o deck):\n"
            f"Tipo de deck: {narrative.get('classification', 'general')}\n"
            f"Audiência: {narrative.get('audience', 'mixed')}\n"
            f"Mensagem central: {narrative.get('central_message', '')}\n"
            f"Arco narrativo: {' → '.join(narrative.get('narrative_arc', []))}\n"
        )
        if narrative.get("key_metrics"):
            narrative_text += "Métricas chave:\n"
            for m in narrative["key_metrics"][:6]:
                narrative_text += f"  - {m.get('label', '')}: {m.get('value', '')} ({m.get('context', '')})\n"
        if narrative.get("comparisons"):
            narrative_text += "Comparações identificadas:\n"
            for c in narrative["comparisons"][:3]:
                narrative_text += f"  - {c.get('entity_a', '')} vs {c.get('entity_b', '')}: {', '.join(c.get('dimensions', [])[:4])}\n"
        if narrative.get("recommendation"):
            narrative_text += f"Recomendação final: {narrative['recommendation']}\n"
        user_prompt_parts.append(narrative_text)

    if benchmark_brief:
        user_prompt_parts.append(
            "BRIEFING EXECUTIVO ESTRUTURADO (usa isto para criar visuais reais, não bullets genéricos):\n"
            f"{benchmark_brief['briefing_text'][:12000]}"
        )
    user_prompt_parts.append(f"CONTEÚDO PARA A APRESENTAÇÃO:\n{planner_content}")

    user_prompt = "\n\n".join(user_prompt_parts)
    planner_prompt = _SLIDE_PLANNER_PROMPT
    if benchmark_brief or executive_mode:
        planner_prompt += "\n\n" + _BENCHMARK_EXEC_APPENDIX
    full_prompt = f"{planner_prompt}\n\n{user_prompt}"

    try:
        raw = await llm_simple(full_prompt, tier=tier, max_tokens=8000)

        slides = _parse_json_from_llm(raw)
        if not isinstance(slides, list) or len(slides) == 0:
            logger.warning("[PptxPlanner] No valid slides from Opus, falling back")
            return _fallback_slides_from_content(content, title)

        # Validate each slide has a type
        validated = []
        for s in slides:
            if isinstance(s, dict) and s.get("type"):
                validated.append(s)
        if not validated:
            return _fallback_slides_from_content(content, title)

        validated = _review_and_rebalance_slide_plan(
            validated,
            title=title,
            content=sanitized_content,
            context=context,
            benchmark_brief=benchmark_brief,
        )

        logger.info("[PptxPlanner] Opus generated %d slides for '%s'",
                    len(validated), (title or content[:50]))
        return validated

    except Exception as e:
        logger.warning("[PptxPlanner] Opus planning failed: %s — falling back", e)
        fallback = _fallback_slides_from_content(content, title)
        return _review_and_rebalance_slide_plan(
            fallback,
            title=title,
            content=sanitized_content,
            context=context,
            benchmark_brief=benchmark_brief,
        )


def _fallback_slides_from_content(content: str, title: str = "") -> List[Dict[str, Any]]:
    """Intelligent fallback: structure content into varied slide types.

    When the AI planner is unavailable, this produces a reasonable presentation
    with section dividers, tables from markdown tables, charts from numeric
    data, max 6 bullets per slide, and unique titles.
    """
    content = _sanitize_planner_content(content)
    lines = [l.strip() for l in content.split("\n") if l.strip()]
    if not lines:
        return [{"type": "content", "title": title or "Conteúdo", "bullets": ["(sem conteúdo)"]}]

    benchmark_brief = _build_benchmark_brief(content, title, "")
    benchmark_sections = _extract_benchmark_sections(content) if benchmark_brief else []
    api_summary = _extract_benchmark_api_summary(content) if benchmark_brief else None
    if benchmark_brief and benchmark_sections:
        slides: List[Dict[str, Any]] = [
            {"type": "agenda", "items": benchmark_brief["agenda_items"], "notes": "Agenda executiva do benchmark."},
            {
                "type": "stat_chart",
                "title": "Síntese executiva",
                "stat_value": benchmark_brief["summary_kpis"][1]["value"],
                "stat_label": "Readiness média estimada",
                "chart_type": "bar",
                "categories": benchmark_brief["chart_categories"],
                "series": [{"name": "Readiness Millennium (%)", "values": benchmark_brief["chart_values"]}],
                "notes": "Abrir com os sinais executivos do benchmark.",
            },
            {
                "type": "chart",
                "title": "Readiness estimada por domínio",
                "chart_type": "bar",
                "categories": benchmark_brief["chart_categories"],
                "series": [{"name": "Readiness Millennium (%)", "values": benchmark_brief["chart_values"]}],
                "notes": "Gap competitivo por domínio.",
            },
            {
                "type": "table",
                "title": "Top gaps com impacto cliente",
                "headers": ["Domínio", "Readiness", "Gap principal", "Impacto"],
                "rows": benchmark_brief["table_rows"],
                "notes": "Tabela de leitura rápida para o comité executivo.",
            },
            {
                "type": "comparison",
                "title": "Posicionamento competitivo",
                "left_title": "Revolut reforça",
                "left_items": benchmark_brief["comparison_left"],
                "right_title": "Millennium precisa reforçar",
                "right_items": benchmark_brief["comparison_right"],
                "notes": "Síntese comparativa do benchmark.",
            },
        ]
        feature_chart_categories = [row["title"] for row in benchmark_brief["domain_rows"][:8]]
        feature_chart_values = [float(row["feature_count"]) for row in benchmark_brief["domain_rows"][:8]]
        slides.append({
            "type": "chart",
            "title": "Amplitude funcional do benchmark por domínio",
            "chart_type": "column",
            "categories": feature_chart_categories,
            "series": [{"name": "N.º de capacidades benchmark", "values": feature_chart_values}],
            "notes": "Mostrar densidade funcional do benchmark em cada domínio.",
        })

        for sec in benchmark_sections[:4]:
            slides.append({
                "type": "comparison",
                "title": sec["title"],
                "left_title": "Benchmark / Revolut",
                "left_items": [item[:90] for item in sec.get("features", [])[:3]] or ["Cobertura benchmark relevante"],
                "right_title": "Gap Millennium",
                "right_items": (
                    [item[:90] for item in sec.get("status", [])[:2]]
                    + [f"Impacto: {item[:82]}" for item in sec.get("impact", [])[:1]]
                )[:3] or ["Gap competitivo material"],
                "notes": f"Detalhar {sec['title']} em modo comparativo executivo.",
            })
        if api_summary:
            slides.extend([
                {
                    "type": "comparison",
                    "title": "Revolut API vs Millennium API",
                    "left_title": "Revolut acelera",
                    "left_items": api_summary["left_items"],
                    "right_title": "Millennium aprofunda",
                    "right_items": api_summary["right_items"],
                    "notes": "Comparar velocidade, profundidade funcional e esforço de integração.",
                },
                {
                    "type": "table",
                    "title": "Leitura executiva das APIs",
                    "headers": ["Tema", "Revolut", "Millennium", "Vantagem"],
                    "rows": api_summary["table_rows"],
                    "notes": "Tabela executiva de apoio para a discussão de APIs.",
                },
            ])
        slides.append({
            "type": "process",
            "title": "Prioridades recomendadas",
            "steps": [
                {"label": "1. Defender base", "description": "Fechar gaps críticos em permissões, pagamentos e controlo."},
                {"label": "2. Criar quick wins", "description": "Atacar capacidades com elevado impacto e esforço moderado."},
                {"label": "3. Reforçar integração", "description": "Clarificar ambição para APIs, ficheiros e ecossistema."},
                {"label": "4. Comunicar valor", "description": "Traduzir o roadmap em proposta executiva clara."},
            ],
            "notes": "Fechar com as prioridades recomendadas para decisão executiva.",
        })
        slides.append({
            "type": "closing",
            "text": "Decisões recomendadas",
            "subtitle": "Priorizar gaps críticos, quick wins e reforço do posicionamento competitivo",
            "notes": "Fechar com o pedido de decisão executiva.",
        })
        return slides[:14]

    # ── Phase 1: Parse into sections with typed content ──
    sections: List[Dict[str, Any]] = []
    current_section: Dict[str, Any] = {"title": title or "Introdução", "bullets": [], "tables": []}
    current_table: Dict[str, Any] = {}  # {headers: [], rows: []}

    def _flush_table():
        nonlocal current_table
        if current_table and current_table.get("headers") and current_table.get("rows"):
            current_section["tables"].append(current_table)
        current_table = {}

    def _flush_section():
        _flush_table()
        if current_section["bullets"] or current_section["tables"]:
            sections.append(current_section)

    for line in lines:
        # Detect headers
        header_match = None
        if line.startswith("# "):
            header_match = line[2:].strip()
        elif line.startswith("## "):
            header_match = line[3:].strip()
        elif line.startswith("### "):
            header_match = line[4:].strip()
        elif re.match(r"^\*\*(.+)\*\*:?\s*$", line):
            header_match = re.match(r"^\*\*(.+)\*\*:?\s*$", line).group(1)

        if header_match:
            _flush_section()
            current_section = {"title": header_match, "bullets": [], "tables": []}
            continue

        # Markdown table separator — skip but signals table mode
        if line.startswith("|") and re.match(r"^\|[\s\-:|]+\|$", line):
            continue

        # Markdown table row
        if line.startswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|") if c.strip()]
            if cells:
                if not current_table.get("headers"):
                    current_table = {"headers": cells, "rows": []}
                else:
                    current_table["rows"].append(cells)
            continue

        # Non-table line → flush any pending table
        _flush_table()

        # Regular bullet
        bullet = line.lstrip("- •*>0123456789.)").strip()
        if bullet and len(bullet) > 2:
            if len(bullet) > 120:
                bullet = bullet[:117] + "..."
            current_section["bullets"].append(bullet)

    _flush_section()

    if not sections:
        return [{"type": "content", "title": title or "Conteúdo", "bullets": lines[:6]}]

    # ── Phase 2: Build slides with variety ──
    slides: List[Dict[str, Any]] = []
    section_num = 0

    for sec in sections:
        section_num += 1
        sec_title = sec["title"]
        bullets = sec["bullets"]
        tables = sec.get("tables", [])

        # Section divider every 2-3 sections
        if len(sections) >= 3 and section_num % 2 == 1:
            slides.append({
                "type": "section",
                "title": sec_title,
                "notes": f"Secção {section_num}: {sec_title}",
            })

        # Emit table slides from detected markdown tables
        for tbl in tables:
            headers = tbl["headers"]
            rows = tbl["rows"]
            if len(rows) > 8:
                # Split large tables across slides
                for chunk_start in range(0, len(rows), 8):
                    chunk_rows = rows[chunk_start:chunk_start + 8]
                    tbl_title = sec_title if chunk_start == 0 else f"{sec_title} (cont.)"
                    slides.append({
                        "type": "table",
                        "title": tbl_title,
                        "headers": headers,
                        "rows": chunk_rows,
                        "notes": f"Tabela: {sec_title} ({len(chunk_rows)} linhas)",
                    })
            else:
                slides.append({
                    "type": "table",
                    "title": sec_title,
                    "headers": headers,
                    "rows": rows,
                    "notes": f"Tabela: {sec_title} ({len(rows)} linhas)",
                })

            # Try to generate a chart from the table data (numeric columns)
            if len(headers) >= 2 and len(rows) >= 2:
                chart_slide = _try_chart_from_table(sec_title, headers, rows)
                if chart_slide:
                    slides.append(chart_slide)

        # Emit content slides (max 6 bullets each)
        for chunk_start in range(0, len(bullets), 6):
            chunk = bullets[chunk_start:chunk_start + 6]
            chunk_title = sec_title
            if chunk_start > 0:
                chunk_title = f"{sec_title} (cont.)"
            slides.append({
                "type": "content",
                "title": chunk_title,
                "bullets": chunk,
                "notes": f"Pontos sobre {sec_title}",
            })

    # Cap at 25 slides max
    slides = slides[:25]

    return slides if slides else [{"type": "content", "title": title, "bullets": lines[:6]}]


def _parse_chart_number(value: str) -> Optional[float]:
    """Parse a string to float, handling European decimal commas and currency/percent symbols.

    Examples:
        "12,5" → 12.5    (European decimal)
        "12,5%" → 12.5   (percent stripped)
        "1.234,56" → 1234.56  (thousands dot + decimal comma)
        "€1200" → 1200.0
        "Sim" → None
    """
    txt = str(value or "").strip()
    if not txt:
        return None
    # Strip currency and percent symbols
    txt = re.sub(r"[€$£%\s]", "", txt)
    if not txt:
        return None
    # Decide decimal/thousands separators by pattern and last separator.
    if "," in txt and "." in txt:
        # Last separator wins as decimal marker:
        # "1.234,56" -> European, "1,234.56" -> English
        if txt.rfind(",") > txt.rfind("."):
            txt = txt.replace(".", "").replace(",", ".")
        else:
            txt = txt.replace(",", "")
    elif "," in txt:
        # "1,234" -> thousands, "12,5" -> decimal
        if re.match(r"^-?\d{1,3}(,\d{3})+$", txt):
            txt = txt.replace(",", "")
        else:
            txt = txt.replace(",", ".")
    elif "." in txt:
        # "1.234" -> thousands, "3.14" -> decimal
        if re.match(r"^-?\d{1,3}(\.\d{3})+$", txt):
            txt = txt.replace(".", "")
    try:
        return float(txt)
    except (ValueError, TypeError):
        return None


def _try_chart_from_table(title: str, headers: List[str], rows: List[List[str]]) -> Optional[Dict[str, Any]]:
    """Try to create a chart slide from table data. Returns slide spec or None.

    Looks for numeric columns to use as chart series, first text column as categories.
    """
    if len(headers) < 2 or len(rows) < 2:
        return None

    # Find label column (first text-like column) and numeric columns
    label_col_idx = None
    numeric_cols: List[int] = []

    for col_idx in range(len(headers)):
        col_values = [row[col_idx] if col_idx < len(row) else "" for row in rows]
        numeric_count = 0
        for val in col_values:
            if _parse_chart_number(str(val)) is not None:
                numeric_count += 1

        if numeric_count >= len(col_values) * 0.6:
            numeric_cols.append(col_idx)
        elif label_col_idx is None:
            label_col_idx = col_idx

    if label_col_idx is None or not numeric_cols:
        return None

    # Max 12 categories, max 3 series for clarity
    categories = [str(row[label_col_idx] if label_col_idx < len(row) else "")
                  for row in rows[:12]]
    series = []
    for col_idx in numeric_cols[:3]:
        values = []
        for row in rows[:12]:
            raw = str(row[col_idx] if col_idx < len(row) else "0")
            parsed = _parse_chart_number(raw)
            values.append(parsed if parsed is not None else 0)
        series.append({
            "name": headers[col_idx] if col_idx < len(headers) else f"Série {col_idx}",
            "values": values,
        })

    chart_type = "bar" if len(series) <= 2 else "column"

    return {
        "type": "chart",
        "title": title,
        "chart_type": chart_type,
        "categories": categories,
        "series": series,
        "notes": f"Gráfico gerado a partir da tabela: {title}",
    }


def _dedupe_slide_titles(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Avoid repeated visible titles in long decks."""
    seen: Dict[str, int] = {}
    deduped = []
    for spec in slides:
        title = str(spec.get("title", "")).strip()
        if not title:
            deduped.append(spec)
            continue
        key = _normalize_text_key(title)
        count = seen.get(key, 0)
        seen[key] = count + 1
        if count > 0 and not title.endswith(")"):
            spec = {**spec, "title": f"{title} ({count + 1})"}
        deduped.append(spec)
    return deduped


def _upgrade_dense_content_slides(slides: List[Dict[str, Any]], *, executive_mode: bool = False) -> List[Dict[str, Any]]:
    """Break monotony by upgrading dense bullet slides into two-column layouts."""
    upgraded = []
    consecutive_content = 0
    for spec in slides:
        slide_type = str(spec.get("type", "")).lower().strip()
        if slide_type == "content":
            consecutive_content += 1
            bullets = list(spec.get("bullets", []) or [])
            if executive_mode and consecutive_content >= 3 and len(bullets) >= 4:
                midpoint = (len(bullets) + 1) // 2
                upgraded.append({
                    "type": "two_column",
                    "title": spec.get("title", ""),
                    "left": bullets[:midpoint],
                    "right": bullets[midpoint:],
                    "notes": spec.get("notes", ""),
                })
                consecutive_content = 0
                continue
        else:
            consecutive_content = 0
        upgraded.append(spec)
    return upgraded


def _upgrade_benchmark_domain_slides(
    slides: List[Dict[str, Any]],
    benchmark_brief: Dict[str, Any],
) -> List[Dict[str, Any]]:
    """Prefer visual comparison slides over plain benchmark content slides."""
    sections = benchmark_brief.get("sections", []) or []
    section_map = {
        _normalize_text_key(sec.get("title", "")): sec
        for sec in sections
        if sec.get("title")
    }
    upgraded = []
    for spec in slides:
        stype = str(spec.get("type", "")).lower().strip()
        title = str(spec.get("title", "")).strip()
        match = section_map.get(_normalize_text_key(title))
        if stype == "content" and match:
            upgraded.append({
                "type": "comparison",
                "title": title,
                "left_title": "Benchmark / Revolut",
                "left_items": match.get("features", [])[:3] or ["Cobertura benchmark relevante"],
                "right_title": "Gap Millennium",
                "right_items": (
                    [item[:90] for item in match.get("status", [])[:2]]
                    + [f"Impacto: {item[:82]}" for item in match.get("impact", [])[:1]]
                )[:3] or [match.get("gap", "Gap competitivo material")],
                "notes": spec.get("notes", ""),
            })
            continue
        if stype == "content" and title:
            norm_title = _normalize_text_key(title)
            if "posicionamento competitivo" in norm_title or "revolut" in norm_title:
                upgraded.append({
                    "type": "comparison",
                    "title": "Posicionamento competitivo",
                    "left_title": "Revolut reforça",
                    "left_items": benchmark_brief.get("comparison_left", [])[:5],
                    "right_title": "Millennium precisa reforçar",
                    "right_items": benchmark_brief.get("comparison_right", [])[:5],
                    "notes": spec.get("notes", ""),
                })
                continue
            if "gap" in norm_title or "millennium" in norm_title:
                upgraded.append({
                    "type": "table",
                    "title": "Top gaps com impacto cliente",
                    "headers": ["Domínio", "Readiness", "Gap principal", "Impacto"],
                    "rows": benchmark_brief.get("table_rows", [])[:5],
                    "notes": spec.get("notes", ""),
                })
                continue
        upgraded.append(spec)
    return upgraded


def _drop_benchmark_decorative_sections(slides: List[Dict[str, Any]], benchmark_brief: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Remove empty-looking section dividers in benchmark decks to keep flow tight."""
    if len(slides) <= 10:
        return slides
    cleaned = []
    for idx, spec in enumerate(slides):
        stype = str(spec.get("type", "")).lower().strip()
        if stype in ("section", "section_divider", "divider"):
            # Keep only if it's the sole separator before a genuinely new chapter;
            # benchmark decks read better when every slide carries information.
            continue
        cleaned.append(spec)
    return cleaned if cleaned else slides


def _inject_benchmark_anchor_slides(
    slides: List[Dict[str, Any]], benchmark_brief: Dict[str, Any], *, executive_mode: bool = False
) -> List[Dict[str, Any]]:
    """Inject mandatory visual/executive anchor slides for benchmark decks."""
    if not benchmark_brief:
        return slides

    result = list(slides)
    types = [str(s.get("type", "")).lower().strip() for s in result]
    chart_count = sum(1 for t in types if t in ("chart", "stat_chart"))

    def _find_insert_idx_after_agenda() -> int:
        for idx, spec in enumerate(result):
            stype = str(spec.get("type", "")).lower().strip()
            if stype == "agenda":
                return idx + 1
        for idx, spec in enumerate(result):
            stype = str(spec.get("type", "")).lower().strip()
            if stype not in ("title", "cover", "capa"):
                return idx
        return len(result)

    if "agenda" not in types and len(result) >= 3:
        agenda_idx = 1 if types and types[0] in ("title", "cover", "capa") else 0
        result.insert(agenda_idx, {
            "type": "agenda",
            "items": benchmark_brief["agenda_items"],
            "notes": "Agenda executiva para enquadrar a comissão executiva.",
        })
        types.insert(agenda_idx, "agenda")

    insert_idx = _find_insert_idx_after_agenda()
    if not any(t in ("kpi", "stat_chart") for t in types):
        result.insert(insert_idx, {
            "type": "kpi",
            "title": "Mensagem executiva",
            "kpis": benchmark_brief["summary_kpis"],
            "notes": "Abrir com os sinais mais fortes do benchmark e orientar a decisão.",
        })
        insert_idx += 1
    if chart_count == 0:
        result.insert(insert_idx, {
            "type": "chart",
            "title": "Readiness estimada por domínio",
            "chart_type": "bar",
            "categories": benchmark_brief["chart_categories"],
            "series": [{"name": "Readiness Millennium (%)", "values": benchmark_brief["chart_values"]}],
            "notes": "Mostrar rapidamente onde o gap competitivo é mais severo.",
        })
        insert_idx += 1
    if chart_count < 2:
        result.insert(insert_idx, {
            "type": "chart",
            "title": "Amplitude funcional do benchmark por domínio",
            "chart_type": "column",
            "categories": [row["title"] for row in benchmark_brief["domain_rows"][:8]],
            "series": [{
                "name": "N.º de capacidades benchmark",
                "values": [float(row["feature_count"]) for row in benchmark_brief["domain_rows"][:8]],
            }],
            "notes": "Mostrar a densidade funcional do benchmark por domínio.",
        })
        insert_idx += 1
    if "table" not in types:
        result.insert(insert_idx, {
            "type": "table",
            "title": "Top gaps com impacto cliente",
            "headers": ["Domínio", "Readiness", "Gap principal", "Impacto"],
            "rows": benchmark_brief["table_rows"],
            "notes": "Tabela de leitura rápida para suportar a discussão executiva.",
        })
        insert_idx += 1
    if "comparison" not in types:
        result.insert(insert_idx, {
            "type": "comparison",
            "title": "Posicionamento competitivo",
            "left_title": "Revolut reforça",
            "left_items": benchmark_brief["comparison_left"],
            "right_title": "Millennium precisa reforçar",
            "right_items": benchmark_brief["comparison_right"],
            "notes": "Comparação resumida entre diferenciação da Revolut e gaps prioritários do Millennium.",
        })
        insert_idx += 1
    if executive_mode and "closing" not in types and "end" not in types and "obrigado" not in types:
        result.append({
            "type": "closing",
            "text": "Decisões recomendadas",
            "subtitle": "Priorizar gaps críticos, quick wins e reforço do posicionamento competitivo",
            "notes": "Fechar com decisão pedida à comissão executiva.",
        })

    return result


def _review_and_rebalance_slide_plan(
    slides: List[Dict[str, Any]],
    *,
    title: str = "",
    content: str = "",
    context: str = "",
    benchmark_brief: Optional[Dict[str, Any]] = None,
) -> List[Dict[str, Any]]:
    """Deterministic quality gate to keep decks executive and visually varied."""
    if not slides:
        return slides

    executive_mode = _looks_like_executive_request(title, context, content) or bool(benchmark_brief)
    rebalanced = list(slides)

    if benchmark_brief:
        rebalanced = _inject_benchmark_anchor_slides(rebalanced, benchmark_brief, executive_mode=executive_mode)

    rebalanced = _upgrade_dense_content_slides(rebalanced, executive_mode=executive_mode)
    if benchmark_brief:
        rebalanced = _upgrade_benchmark_domain_slides(rebalanced, benchmark_brief)
        rebalanced = _drop_benchmark_decorative_sections(rebalanced, benchmark_brief)
    rebalanced = _dedupe_slide_titles(rebalanced)

    # Keep the deck focused for executive audiences.
    if executive_mode and len(rebalanced) > 14:
        protected_types = {"title", "cover", "capa", "agenda", "kpi", "stat_chart", "chart", "table", "comparison", "closing", "end", "obrigado"}
        trimmed: List[Dict[str, Any]] = []
        content_budget = 14
        for spec in rebalanced:
            stype = str(spec.get("type", "")).lower().strip()
            if stype in protected_types:
                trimmed.append(spec)
                continue
            if len(trimmed) < content_budget:
                trimmed.append(spec)
        rebalanced = trimmed

    return rebalanced
